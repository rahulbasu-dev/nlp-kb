"""Generate detailed educational PPTX presentations for 4 NLP topics.

Topics:
  1. Word2Vec
  2. SGNS (Skip-gram with Negative Sampling)
  3. CBOW (Continuous Bag of Words)
  4. GloVe (Global Vectors for Word Representation)

All slides use white backgrounds with clean, professional styling.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE


# ── Color palette (white-background theme) ────────────────────────────────
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x1A, 0x1A, 0x2E)
DARK_GRAY = RGBColor(0x33, 0x33, 0x33)
MED_GRAY = RGBColor(0x66, 0x66, 0x66)
LIGHT_GRAY = RGBColor(0xE8, 0xE8, 0xE8)
ACCENT_BLUE = RGBColor(0x1A, 0x73, 0xE8)
ACCENT_DARK_BLUE = RGBColor(0x0D, 0x47, 0xA1)
ACCENT_GREEN = RGBColor(0x0F, 0x9D, 0x58)
ACCENT_ORANGE = RGBColor(0xF4, 0xB4, 0x00)
ACCENT_RED = RGBColor(0xDB, 0x44, 0x37)
BG_LIGHT_BLUE = RGBColor(0xE8, 0xF0, 0xFE)
BG_LIGHT_GREEN = RGBColor(0xE6, 0xF4, 0xEA)
BG_LIGHT_ORANGE = RGBColor(0xFE, 0xF7, 0xE0)
BG_LIGHT_GRAY = RGBColor(0xF5, 0xF5, 0xF5)


def _set_slide_bg(slide, color: RGBColor = WHITE) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def _add_rect(slide, left, top, width, height, fill_color,
              line_color=None, line_width=Pt(0)):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.fill.solid()
        shape.line.fill.fore_color.rgb = line_color
        shape.line.width = line_width
    else:
        shape.line.fill.background()
    return shape


def _add_rounded_rect(slide, left, top, width, height, fill_color,
                      line_color=None, line_width=Pt(1)):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.fill.solid()
        shape.line.fill.fore_color.rgb = line_color
        shape.line.width = line_width
    else:
        shape.line.fill.background()
    return shape


def _add_textbox(slide, left, top, width, height):
    return slide.shapes.add_textbox(left, top, width, height)


def _set_para(para, text, size=14, color=DARK_GRAY, bold=False,
              italic=False, alignment=PP_ALIGN.LEFT):
    para.text = text
    para.font.size = Pt(size)
    para.font.color.rgb = color
    para.font.bold = bold
    para.font.italic = italic
    para.alignment = alignment


def _add_para(tf, text, size=14, color=DARK_GRAY, bold=False,
              italic=False, alignment=PP_ALIGN.LEFT, space_before=Pt(4),
              space_after=Pt(2)):
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.italic = italic
    p.alignment = alignment
    p.space_before = space_before
    p.space_after = space_after
    return p


def _new_prs():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    return prs


def _blank_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, WHITE)
    return slide


# ── Reusable slide builders ───────────────────────────────────────────────

def add_title_slide(prs, title: str, subtitle: str, accent=ACCENT_BLUE):
    """Full-width title slide with accent bar."""
    slide = _blank_slide(prs)
    # Accent bar at top
    _add_rect(slide, Inches(0), Inches(0),
              prs.slide_width, Inches(0.12), accent)
    # Title
    tb = _add_textbox(slide, Inches(1), Inches(2.0), Inches(11.3), Inches(1.8))
    tf = tb.text_frame
    tf.word_wrap = True
    _set_para(tf.paragraphs[0], title, 40, BLACK, True, alignment=PP_ALIGN.CENTER)
    # Subtitle
    tb2 = _add_textbox(slide, Inches(2), Inches(4.0), Inches(9.3), Inches(1))
    tf2 = tb2.text_frame
    tf2.word_wrap = True
    _set_para(tf2.paragraphs[0], subtitle, 20, MED_GRAY, False,
              alignment=PP_ALIGN.CENTER)
    # Thin line
    _add_rect(slide, Inches(5), Inches(5.3), Inches(3.3), Pt(2), accent)
    return slide


def add_section_slide(prs, section_title: str, accent=ACCENT_BLUE):
    """Section divider slide."""
    slide = _blank_slide(prs)
    _add_rect(slide, Inches(0), Inches(0),
              prs.slide_width, Inches(0.08), accent)
    _add_rect(slide, Inches(1), Inches(3.2), Inches(0.15), Inches(1.2), accent)
    tb = _add_textbox(slide, Inches(1.5), Inches(3.2), Inches(10), Inches(1.2))
    tf = tb.text_frame
    tf.word_wrap = True
    _set_para(tf.paragraphs[0], section_title, 34, BLACK, True)
    return slide


def add_content_slide(prs, title: str, bullets: list,
                      accent=ACCENT_BLUE, subtitle: str = ""):
    """Standard content slide with title and bullet points."""
    slide = _blank_slide(prs)
    # Accent bar
    _add_rect(slide, Inches(0), Inches(0),
              prs.slide_width, Inches(0.06), accent)
    # Title
    tb = _add_textbox(slide, Inches(0.6), Inches(0.3), Inches(12), Inches(0.8))
    tf = tb.text_frame
    tf.word_wrap = True
    _set_para(tf.paragraphs[0], title, 28, BLACK, True)
    if subtitle:
        _add_para(tf, subtitle, 16, MED_GRAY, italic=True)
    # Bullets
    top = 1.4 if not subtitle else 1.7
    tb2 = _add_textbox(slide, Inches(0.8), Inches(top),
                       Inches(11.5), Inches(7.5 - top - 0.3))
    tf2 = tb2.text_frame
    tf2.word_wrap = True
    for i, bullet in enumerate(bullets):
        if i == 0:
            p = tf2.paragraphs[0]
        else:
            p = tf2.add_paragraph()
        # Support indented sub-bullets with "  - " prefix
        if bullet.startswith("  - "):
            p.text = bullet.strip("  - ").strip()
            p.level = 1
            p.font.size = Pt(15)
        else:
            p.text = bullet
            p.font.size = Pt(17)
        p.font.color.rgb = DARK_GRAY
        p.space_before = Pt(6)
        p.space_after = Pt(3)
    return slide


def add_two_column_slide(prs, title: str, left_title: str, left_bullets: list,
                         right_title: str, right_bullets: list,
                         accent=ACCENT_BLUE):
    """Two-column comparison slide."""
    slide = _blank_slide(prs)
    _add_rect(slide, Inches(0), Inches(0),
              prs.slide_width, Inches(0.06), accent)
    # Title
    tb = _add_textbox(slide, Inches(0.6), Inches(0.3), Inches(12), Inches(0.7))
    _set_para(tb.text_frame.paragraphs[0], title, 28, BLACK, True)
    # Left column card
    card_l = _add_rounded_rect(slide, Inches(0.5), Inches(1.3),
                               Inches(5.9), Inches(5.7), BG_LIGHT_BLUE,
                               ACCENT_BLUE, Pt(1))
    tb_lt = _add_textbox(slide, Inches(0.8), Inches(1.5),
                         Inches(5.3), Inches(0.5))
    _set_para(tb_lt.text_frame.paragraphs[0], left_title,
              20, ACCENT_DARK_BLUE, True)
    tb_lb = _add_textbox(slide, Inches(0.8), Inches(2.2),
                         Inches(5.3), Inches(4.5))
    tf_lb = tb_lb.text_frame
    tf_lb.word_wrap = True
    for i, b in enumerate(left_bullets):
        p = tf_lb.paragraphs[0] if i == 0 else tf_lb.add_paragraph()
        p.text = b
        p.font.size = Pt(15)
        p.font.color.rgb = DARK_GRAY
        p.space_before = Pt(5)
        p.space_after = Pt(2)
    # Right column card
    card_r = _add_rounded_rect(slide, Inches(6.9), Inches(1.3),
                               Inches(5.9), Inches(5.7), BG_LIGHT_GREEN,
                               ACCENT_GREEN, Pt(1))
    tb_rt = _add_textbox(slide, Inches(7.2), Inches(1.5),
                         Inches(5.3), Inches(0.5))
    _set_para(tb_rt.text_frame.paragraphs[0], right_title,
              20, ACCENT_GREEN, True)
    tb_rb = _add_textbox(slide, Inches(7.2), Inches(2.2),
                         Inches(5.3), Inches(4.5))
    tf_rb = tb_rb.text_frame
    tf_rb.word_wrap = True
    for i, b in enumerate(right_bullets):
        p = tf_rb.paragraphs[0] if i == 0 else tf_rb.add_paragraph()
        p.text = b
        p.font.size = Pt(15)
        p.font.color.rgb = DARK_GRAY
        p.space_before = Pt(5)
        p.space_after = Pt(2)
    return slide


def add_formula_slide(prs, title: str, formulas: list, explanations: list,
                      accent=ACCENT_BLUE):
    """Slide highlighting formulas with explanations."""
    slide = _blank_slide(prs)
    _add_rect(slide, Inches(0), Inches(0),
              prs.slide_width, Inches(0.06), accent)
    tb = _add_textbox(slide, Inches(0.6), Inches(0.3), Inches(12), Inches(0.7))
    _set_para(tb.text_frame.paragraphs[0], title, 28, BLACK, True)

    top = 1.4
    for formula, explanation in zip(formulas, explanations):
        # Formula box
        card = _add_rounded_rect(slide, Inches(0.8), Inches(top),
                                 Inches(11.7), Inches(0.7), BG_LIGHT_GRAY,
                                 LIGHT_GRAY, Pt(1))
        ftb = _add_textbox(slide, Inches(1.0), Inches(top + 0.1),
                           Inches(11.3), Inches(0.5))
        _set_para(ftb.text_frame.paragraphs[0], formula,
                  18, ACCENT_DARK_BLUE, True, alignment=PP_ALIGN.CENTER)
        top += 0.8
        # Explanation
        etb = _add_textbox(slide, Inches(1.2), Inches(top),
                           Inches(11), Inches(0.6))
        etf = etb.text_frame
        etf.word_wrap = True
        _set_para(etf.paragraphs[0], explanation, 15, MED_GRAY)
        top += 0.7
    return slide


def add_table_slide(prs, title: str, headers: list, rows: list,
                    accent=ACCENT_BLUE, col_widths=None):
    """Slide with a styled table."""
    slide = _blank_slide(prs)
    _add_rect(slide, Inches(0), Inches(0),
              prs.slide_width, Inches(0.06), accent)
    tb = _add_textbox(slide, Inches(0.6), Inches(0.3), Inches(12), Inches(0.7))
    _set_para(tb.text_frame.paragraphs[0], title, 28, BLACK, True)

    n_cols = len(headers)
    n_rows = len(rows) + 1
    table_width = Inches(11.5)
    table_left = Inches(0.9)

    tbl_shape = slide.shapes.add_table(
        n_rows, n_cols, table_left, Inches(1.4), table_width, Inches(0.5 * n_rows))
    tbl = tbl_shape.table

    # Set column widths
    if col_widths:
        for i, w in enumerate(col_widths):
            tbl.columns[i].width = Inches(w)

    # Header row
    for j, h in enumerate(headers):
        cell = tbl.cell(0, j)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = accent
        for para in cell.text_frame.paragraphs:
            para.font.size = Pt(14)
            para.font.color.rgb = WHITE
            para.font.bold = True
            para.alignment = PP_ALIGN.CENTER

    # Data rows
    for i, row in enumerate(rows):
        for j, val in enumerate(row):
            cell = tbl.cell(i + 1, j)
            cell.text = str(val)
            cell.fill.solid()
            cell.fill.fore_color.rgb = WHITE if i % 2 == 0 else BG_LIGHT_GRAY
            for para in cell.text_frame.paragraphs:
                para.font.size = Pt(13)
                para.font.color.rgb = DARK_GRAY
                para.alignment = PP_ALIGN.CENTER
    return slide


def add_key_concept_slide(prs, title: str, concept: str, details: list,
                          accent=ACCENT_BLUE):
    """Highlighted key concept with supporting details."""
    slide = _blank_slide(prs)
    _add_rect(slide, Inches(0), Inches(0),
              prs.slide_width, Inches(0.06), accent)
    tb = _add_textbox(slide, Inches(0.6), Inches(0.3), Inches(12), Inches(0.7))
    _set_para(tb.text_frame.paragraphs[0], title, 28, BLACK, True)
    # Concept highlight box
    card = _add_rounded_rect(slide, Inches(1.5), Inches(1.5),
                             Inches(10.3), Inches(1.2), BG_LIGHT_BLUE,
                             ACCENT_BLUE, Pt(2))
    ctb = _add_textbox(slide, Inches(1.8), Inches(1.65),
                       Inches(9.7), Inches(1.0))
    ctf = ctb.text_frame
    ctf.word_wrap = True
    _set_para(ctf.paragraphs[0], concept, 22, ACCENT_DARK_BLUE, True,
              alignment=PP_ALIGN.CENTER)
    # Details
    dtb = _add_textbox(slide, Inches(1.0), Inches(3.2),
                       Inches(11.3), Inches(4))
    dtf = dtb.text_frame
    dtf.word_wrap = True
    for i, d in enumerate(details):
        p = dtf.paragraphs[0] if i == 0 else dtf.add_paragraph()
        p.text = d
        p.font.size = Pt(16)
        p.font.color.rgb = DARK_GRAY
        p.space_before = Pt(6)
        p.space_after = Pt(3)
    return slide


# ══════════════════════════════════════════════════════════════════════════
# TOPIC 1: WORD2VEC
# ══════════════════════════════════════════════════════════════════════════

def build_word2vec(prs):
    add_title_slide(prs, "Word2Vec", "Neural Word Embeddings\nFrom Words to Vectors")

    # ── Section 1: Intuition ──
    add_section_slide(prs, "1. Intuition")

    add_key_concept_slide(prs, "What is Word2Vec?",
        '"You shall know a word by the company it keeps." -- J.R. Firth (1957)',
        [
            "Word2Vec is a family of neural network models that learn to represent words "
            "as dense vectors in a continuous vector space.",
            "Unlike sparse representations (one-hot, TF-IDF) where most values are zero, "
            "Word2Vec creates compact vectors where every dimension carries meaning.",
            "Words that appear in similar contexts (e.g., 'dog' and 'cat' near 'pet', "
            "'food', 'vet') get similar vector representations.",
            "The result is a 'map' of language where similar words are neighbors, "
            "relationships are preserved, and directions encode meaning.",
        ])

    add_two_column_slide(prs, "Two Architectures: Skip-gram vs. CBOW",
        "Skip-gram",
        [
            '"Given a word, predict its context."',
            "",
            'Example: "The quick brown fox jumps"',
            "  Center: fox (window=2)",
            '  Predict: "quick", "brown", "jumps", "over"',
            "",
            "Best for: Rare words, smaller datasets",
            "Generates more training pairs per position",
        ],
        "CBOW (Continuous Bag of Words)",
        [
            '"Given the context, predict the word."',
            "",
            'Example: "The quick brown ___ jumps"',
            '  Context: "quick", "brown", "jumps", "over"',
            '  Predict: "fox"',
            "",
            "Best for: Frequent words, larger datasets",
            "Faster training, averages context vectors",
        ])

    add_table_slide(prs, "Why Dense Vectors?  Sparse vs. Dense Representations",
        ["Property", "Sparse (One-Hot / TF-IDF)", "Dense (Word2Vec)"],
        [
            ["Vector Size", "Vocabulary size (10,000+)", "Fixed small (100-300)"],
            ["Most Values", "Zero", "Non-zero"],
            ["Similarity", '"cat" and "dog" are orthogonal', '"cat" and "dog" are close'],
            ["Arithmetic", "Meaningless", "Captures relationships!"],
            ["Storage", "Sparse matrix tricks needed", "Compact"],
        ])

    # ── Section 2: Worked Examples ──
    add_section_slide(prs, "2. Worked Examples")

    add_key_concept_slide(prs, "The Famous Analogy: King - Man + Woman = Queen",
        "vec(king) - vec(man) + vec(woman)  ~=  vec(queen)",
        [
            "By subtracting 'man' from 'king', we remove the 'maleness' component, "
            "leaving 'royalty'. Adding 'woman' gives 'female royalty' = 'queen'.",
            "This works because Word2Vec encodes semantic relationships as consistent "
            "vector offsets in the embedding space.",
            "Limitation: Analogies are not perfect. They work best for common, "
            "unambiguous relationships. The model can also encode biases from training data.",
        ])

    add_table_slide(prs, "Analogy Examples",
        ["Relationship", "A", "is to B", "as C", "is to D (predicted)"],
        [
            ["Gender", "man", "woman", "king", "queen"],
            ["Country-Capital", "France", "Paris", "Japan", "Tokyo"],
            ["Verb Tense", "walk", "walked", "swim", "swam"],
            ["Comparative", "big", "bigger", "small", "smaller"],
            ["Superlative", "good", "best", "bad", "worst"],
            ["Profession", "doctor", "hospital", "teacher", "school"],
        ])

    add_table_slide(prs, "Word Similarity via Cosine Similarity",
        ["Word Pair", "Cosine Similarity", "Interpretation"],
        [
            ["cat - kitten", "0.89", "Very similar (same animal)"],
            ["cat - dog", "0.76", "Similar (both pets)"],
            ["cat - car", "0.21", "Unrelated"],
            ["good - bad", "0.72", "Related (antonyms share context)"],
            ["king - queen", "0.85", "Very similar (both royalty)"],
            ["Paris - France", "0.78", "Related (capital-country)"],
        ])

    add_content_slide(prs, "Why Are Antonyms Similar?", [
        '"Good" and "bad" have high similarity because they appear in similar contexts:',
        '  - "The movie was good/bad"',
        '  - "a good/bad idea"',
        "Word2Vec captures relatedness, not just synonymy.",
        "This is actually useful -- it groups semantically related concepts together.",
    ])

    # ── Section 3: Mathematical Foundations ──
    add_section_slide(prs, "3. Mathematical Foundations")

    add_formula_slide(prs, "Skip-gram Objective Function",
        [
            "J = (1/T) * SUM_t SUM_{-c <= j <= c, j != 0}  log P(w_{t+j} | w_t)",
        ],
        [
            "T = total words in corpus;  c = context window size;  "
            "w_t = center word at position t;  w_{t+j} = context word at offset j.",
        ])

    add_formula_slide(prs, "Softmax Probability",
        [
            "P(w_O | w_I) = exp(v'_{w_O} . v_{w_I}) / SUM_{w=1}^{V} exp(v'_w . v_{w_I})",
        ],
        [
            "v_{w_I} = input (center word) embedding;  v'_{w_O} = output (context word) "
            "embedding;  V = vocabulary size.  The denominator sums over ALL words -- "
            "this is the computational bottleneck (100K+ words).",
        ])

    add_content_slide(prs, "The Softmax Problem and Solutions", [
        "Computing full softmax requires summing over the entire vocabulary for every "
        "training example -- computationally prohibitive!",
        "",
        "Solution 1: Hierarchical Softmax",
        "  - Organize vocabulary in a binary tree",
        "  - Traverse from root to target word, making binary decisions at each node",
        "  - Complexity: O(V) --> O(log2 V)  -- for V=100K, this is ~6,000x speedup",
        "",
        "Solution 2: Negative Sampling (most popular)",
        "  - Instead of all words, only update the correct word + small sample of "
        "'negative' (incorrect) words",
        "  - Transforms from multiclass to binary classification",
        "  - Typically k = 5-20 negative samples per positive pair",
    ])

    add_formula_slide(prs, "Negative Sampling Objective",
        [
            "J = log sigma(v'_{w_O} . v_{w_I}) + SUM_{i=1}^{k} E_{w_i ~ P_n(w)} "
            "[log sigma(-v'_{w_i} . v_{w_I})]",
        ],
        [
            "sigma(x) = 1/(1 + e^{-x}) is the sigmoid;  k = negative samples (5-20);  "
            "P_n(w) = noise distribution, usually f(w)^{3/4} / Z.  "
            "The model learns to distinguish real context pairs from random noise pairs.",
        ])

    add_formula_slide(prs, "Cosine Similarity",
        [
            "cos(theta) = (A . B) / (||A|| * ||B||) = SUM(A_i * B_i) / "
            "(sqrt(SUM A_i^2) * sqrt(SUM B_i^2))",
        ],
        [
            "Ranges from -1 (opposite) to 1 (identical);  0 = orthogonal/unrelated.  "
            "Independent of vector magnitude -- focuses on direction only.",
        ])

    # ── Section 4: Numerical Example ──
    add_section_slide(prs, "4. Numerical Example")

    add_content_slide(prs, "King-Queen Analogy with Real Numbers (3D Simplified)", [
        "Imagine 3 dimensions encode: [royalty, gender (male=1), age]",
        "",
        "king  = [1, 1, 1]   (royal, male, adult)",
        "man   = [0, 1, 1]   (not royal, male, adult)",
        "woman = [0, 0, 1]   (not royal, female, adult)",
        "queen = [1, 0, 1]   (royal, female, adult)",
    ])

    add_content_slide(prs, "Step-by-Step Vector Arithmetic", [
        "Step 1:  king - man",
        "  - [1, 1, 1] - [0, 1, 1] = [1, 0, 0]",
        '  - This extracts the "royalty" concept',
        "",
        "Step 2:  (king - man) + woman",
        "  - [1, 0, 0] + [0, 0, 1] = [1, 0, 1]",
        '  - Add "femaleness" to royalty',
        "",
        "Result: [1, 0, 1] = queen!",
        "",
        "In real Word2Vec (100-300 dimensions), the result would be close to "
        "(not exactly) the queen vector. We use cosine similarity to find the nearest word.",
    ])

    add_content_slide(prs,
        "Cosine Similarity Calculation: king [1,1,1] vs queen [1,0,1]", [
        "Step 1: Dot product  A . B = (1)(1) + (1)(0) + (1)(1) = 2",
        "Step 2: ||A|| = sqrt(1 + 1 + 1) = sqrt(3) = 1.732",
        "Step 3: ||B|| = sqrt(1 + 0 + 1) = sqrt(2) = 1.414",
        "Step 4: cos(theta) = 2 / (1.732 * 1.414) = 2 / 2.449 = 0.816",
        "",
        "Interpretation: 0.816 indicates king and queen are highly similar -- "
        "they share the 'royalty' and 'adult' dimensions.",
    ])

    # ── Section 5: Training Pipeline ──
    add_section_slide(prs, "5. Training Pipeline")

    add_content_slide(prs, "Word2Vec Training Pipeline (4 Steps)", [
        "Step 1: Preprocessing",
        "  - Tokenize text into words",
        "  - Build vocabulary (word-to-index mapping)",
        "  - Subsample frequent words (randomly remove 'the', 'is', etc.)",
        "",
        "Step 2: Generate Training Pairs",
        "  - Slide a window across text, create (center, context) pairs",
        "  - Window size c typically 2-10 words on each side",
        "",
        "Step 3: Neural Network Forward Pass",
        "  - Skip-gram: Look up center word embedding, predict context probabilities",
        "  - CBOW: Average context embeddings, predict center word probability",
        "",
        "Step 4: Loss and Backpropagation",
        "  - Compute loss using negative sampling (most common) or hierarchical softmax",
        "  - Update embeddings via gradient descent",
    ])

    # ── Section 6: Applications & Alternatives ──
    add_section_slide(prs, "6. Applications & When to Use")

    add_content_slide(prs, "Applications of Word2Vec", [
        "1. Semantic Search",
        '  - Find documents with synonyms/related terms, not just exact keywords',
        '  - Search "automobile" and find documents about "car"',
        "",
        "2. Recommendation Systems",
        '  - Treat products as "words" and sessions as "sentences" to learn item embeddings',
        "",
        "3. Named Entity Recognition",
        "  - Use embeddings as features for sequence labeling; similar syntactic roles cluster",
        "",
        "4. Machine Translation",
        '  - Embed words from different languages into the same space',
        '  - "roi" (French) and "king" (English) become neighbors',
    ])

    add_table_slide(prs, "When to Use Word2Vec vs. Modern Alternatives",
        ["Method", "Strengths", "Best For"],
        [
            ["Word2Vec", "Fast, interpretable, good baseline", "General tasks"],
            ["FastText", "Handles OOV words via subwords",
             "Morphologically rich languages"],
            ["GloVe", "Combines count-based & neural", "Similar to Word2Vec"],
            ["BERT / GPT", "Contextual (different vectors by context)",
             "State-of-the-art, heavier"],
        ])

    # ── References ──
    add_section_slide(prs, "7. References")
    add_content_slide(prs, "Key References", [
        "[1] Mikolov et al. (2013) - 'Efficient estimation of word representations "
        "in vector space' -- Original Word2Vec paper (Skip-gram & CBOW)",
        "",
        "[2] Mikolov et al. (2013) - 'Distributed representations of words and phrases "
        "and their compositionality' -- Negative sampling, phrase embeddings, king-queen analogy",
        "",
        "[3] Goldberg & Levy (2014) - 'word2vec Explained' -- Clear mathematical derivation "
        "of the negative sampling objective",
        "",
        "[4] Rong (2014) - 'word2vec Parameter Learning Explained' -- Detailed tutorial "
        "on backpropagation and gradient updates",
        "",
        "[5] Levy & Goldberg (2014) - 'Neural Word Embedding as Implicit Matrix "
        "Factorization' -- Shows SGNS implicitly factorizes a shifted PMI matrix",
    ])


# ══════════════════════════════════════════════════════════════════════════
# TOPIC 2: SGNS
# ══════════════════════════════════════════════════════════════════════════

def build_sgns(prs):
    add_title_slide(prs, "Skip-gram with Negative Sampling\n(SGNS)",
                    "Neural Word Embeddings Made Efficient",
                    ACCENT_GREEN)

    # ── Section 1: Intuition ──
    add_section_slide(prs, "1. Intuition", ACCENT_GREEN)

    add_key_concept_slide(prs, "What is Skip-gram?",
        "Given a word, predict the words that appear nearby.",
        [
            "A neural network architecture that learns word representations by training "
            "on a simple task: predict context words from a center word.",
            "Key insight: Words appearing in similar contexts have similar meanings "
            "(distributional hypothesis).",
            'Even without knowing meanings, "dog" and "cat" appear in similar contexts '
            '("adopted a ___", "___ loves to play") -- so they are related.',
        ], ACCENT_GREEN)

    add_content_slide(prs, "Skip-gram: Predicting Context", [
        'Sentence: "The quick brown fox jumps"',
        'Center word: "brown" (window size = 2)',
        'Task: Predict "The", "quick", "fox", "jumps"',
        "This generates 4 training examples from one center word!",
        "",
        'Why "Skip-gram"?  The model can "skip" over words -- context words',
        "don't need to be immediately adjacent, just within the window.",
    ], ACCENT_GREEN)

    add_content_slide(prs, "The Softmax Problem", [
        "The naive approach uses softmax over the entire vocabulary:",
        "  P(context | center) = softmax(scores)",
        "",
        "With a vocabulary of 1 million words, each training step requires:",
        "  - 1M dot products",
        "  - Normalization over 1M probabilities",
        "  - Gradient updates for 1M words",
        "",
        "Result: Impossibly slow! Training would take years.",
    ], ACCENT_GREEN)

    add_key_concept_slide(prs, "The Negative Sampling Breakthrough",
        "Transform multiclass classification (which of V words?) "
        "into binary classification (is this pair real or fake?)",
        [
            "Instead of: 'Which of these 1 million words is the answer?' (multiple choice)",
            "We ask: 'Is this word pair valid?' (true/false)",
            "",
            'Positive pair: ("brown", "fox") --> Yes, these appear together!',
            'Negative pair: ("brown", "refrigerator") --> No, randomly sampled noise!',
            "",
            "Only requires updating a few words per step, not millions!",
        ], ACCENT_GREEN)

    add_content_slide(prs, "Two Embedding Matrices & Subsampling", [
        "Two Embedding Matrices:",
        "  - W: Center word embeddings",
        "  - C: Context word embeddings",
        "  - Final embeddings typically use W or W+C",
        "  - Levy & Goldberg showed SGNS implicitly factorizes a shifted PMI matrix",
        "",
        "Subsampling Frequent Words:",
        "  - Common words ('the', 'is', 'and') provide little semantic information",
        "  - P(keep) = sqrt(t / f(w)) + t / f(w)   where t ~ 10^-5",
        "  - Rare words: almost always kept;  Common words: often discarded",
        "  - Benefits: reduces training time, increases effective window, improves quality",
    ], ACCENT_GREEN)

    # ── Section 2: Worked Examples ──
    add_section_slide(prs, "2. Worked Examples", ACCENT_GREEN)

    add_content_slide(prs, "Example 1: Extracting Training Pairs", [
        'Sentence: "the king loved the queen and the people"',
        "Window size = 2, center word = 'loved'",
        "",
        "Window: [the, king, loved, the, queen]",
        "Positive pairs: (loved, the), (loved, king), (loved, the), (loved, queen)",
        "",
        "For each positive pair, sample k=3 negative words:",
        "  (loved, king)  --> Neg: computer, banana, zebra",
        "  (loved, queen) --> Neg: table, river, planet",
        "  (loved, the)   --> Neg: guitar, mountain, coffee",
    ], ACCENT_GREEN)

    add_content_slide(prs, 'Example 2: Training Step for ("loved", "queen")', [
        "Negative samples: [computer, banana]",
        "",
        "Step 1 -- Look up embeddings (dim=4):",
        "  w_loved = [0.2, -0.5, 0.1, 0.8]",
        "  c_queen = [0.3, 0.4, -0.2, 0.1]",
        "  c_computer = [-0.1, 0.2, 0.5, -0.3]",
        "  c_banana = [0.4, -0.1, 0.3, 0.2]",
        "",
        "Step 2 -- Dot products:",
        "  score_queen = 0.2*0.3 + (-0.5)*0.4 + 0.1*(-0.2) + 0.8*0.1 = -0.08",
        "  score_computer = -0.37",
        "  score_banana = 0.32",
    ], ACCENT_GREEN)

    add_content_slide(prs, "Example 2 (continued): Sigmoid & Gradient Updates", [
        "Step 3 -- Apply sigmoid:  sigma(x) = 1 / (1 + exp(-x))",
        "  P(real | loved, queen) = sigma(-0.08) = 0.48",
        "  P(real | loved, computer) = sigma(-0.37) = 0.41",
        "  P(real | loved, banana) = sigma(0.32) = 0.58",
        "",
        "Step 4 -- Errors (target: queen=1, computer=0, banana=0):",
        "  queen:    1 - 0.48 = 0.52  (push closer)",
        "  computer: 0 - 0.41 = -0.41  (push away)",
        "  banana:   0 - 0.58 = -0.58  (push away)",
        "",
        "Each update adjusts vectors: positive pairs move closer, negative pairs apart.",
        "After millions of updates, co-occurring words cluster together!",
    ], ACCENT_GREEN)

    # ── Section 3: Tiny Corpus Numerical Example ──
    add_section_slide(prs, "3. Full Numerical Example: 'cat sat mat'", ACCENT_GREEN)

    add_content_slide(prs, 'Tiny Corpus: "cat sat mat" (V=3, window=1)', [
        "Training pairs generated:",
        "  Position 0 (cat): cat --> sat",
        "  Position 1 (sat): sat --> cat,  sat --> mat",
        "  Position 2 (mat): mat --> sat",
        "  Total: 4 positive pairs",
        "",
        "Center Word Matrix W (random init, dim=2):",
        "  cat = [0.5, -0.3],  sat = [-0.2, 0.4],  mat = [0.1, 0.6]",
        "",
        "Context Word Matrix C (random init, dim=2):",
        "  cat = [0.3, 0.2],  sat = [-0.4, 0.1],  mat = [0.2, -0.5]",
    ], ACCENT_GREEN)

    add_content_slide(prs,
        'Training on ("cat", "sat") with negative sample "mat"', [
        "Step 1 -- Vectors:  w_cat=[0.5,-0.3], c_sat=[-0.4,0.1], c_mat=[0.2,-0.5]",
        "",
        "Step 2 -- Dot products:",
        "  Positive (cat,sat): 0.5*(-0.4) + (-0.3)*0.1 = -0.23",
        "  Negative (cat,mat): 0.5*0.2 + (-0.3)*(-0.5) = 0.25",
        "",
        "Step 3 -- Sigmoid:",
        "  P(real|cat,sat) = sigma(-0.23) = 0.443",
        "  P(real|cat,mat) = sigma(0.25) = 0.562",
        "",
        "Problem! Model thinks (cat,mat) is MORE likely real than (cat,sat). Backwards!",
    ], ACCENT_GREEN)

    add_content_slide(prs, "Gradient Updates (learning rate = 0.1)", [
        "Errors:  positive = 0.443 - 1 = -0.557;  negative = 0.562 - 0 = 0.562",
        "",
        "Update w_cat from positive: grad = -0.557 * [-0.4, 0.1] = [0.223, -0.056]",
        "  w_cat = [0.5, -0.3] - 0.1*[0.223, -0.056] = [0.478, -0.294]",
        "",
        "Update w_cat from negative: grad = 0.562 * [0.2, -0.5] = [0.112, -0.281]",
        "  w_cat = [0.478, -0.294] - 0.1*[0.112, -0.281] = [0.467, -0.266]",
        "",
        "Before:  dot(cat,sat)=-0.23, dot(cat,mat)=0.25",
        "After:   dot(cat,sat)=-0.20 (improved!), dot(cat,mat)=0.21 (improved!)",
        "",
        "After thousands of such updates, real co-occurring words will be similar!",
    ], ACCENT_GREEN)

    # ── Section 4: Mathematical Foundations ──
    add_section_slide(prs, "4. Mathematical Foundations", ACCENT_GREEN)

    add_formula_slide(prs, "Skip-gram Objective (Full Softmax)",
        [
            "J = (1/T) SUM_{t=1}^{T} SUM_{-c<=j<=c, j!=0}  log P(w_{t+j} | w_t)",
            "P(o|c) = exp(u_o^T * v_c) / SUM_{w=1}^{V} exp(u_w^T * v_c)",
        ],
        [
            "T = total words;  c = window size;  v_c = center embedding (W);  "
            "u_o = context embedding (C).",
            "The denominator sums over ALL V words -- this is the expensive part!",
        ], ACCENT_GREEN)

    add_formula_slide(prs, "Negative Sampling Objective",
        [
            "J = log sigma(u_o^T v_c) + SUM_{i=1}^{k} E_{w_i~P_n} "
            "[log sigma(-u_{w_i}^T v_c)]",
        ],
        [
            "First term: maximize dot product for positive pairs (push together).  "
            "Second term: minimize dot product for negative pairs (push apart).  "
            "k = 5-20 for small datasets, 2-5 for large.  "
            "P_n(w) proportional to f(w)^{0.75} -- the 0.75 exponent is the sweet spot.",
        ], ACCENT_GREEN)

    add_formula_slide(prs, "Gradient Update Rules",
        [
            "dJ/dv_c = (sigma(u_o^T v_c) - 1) u_o + SUM_i sigma(u_{w_i}^T v_c) u_{w_i}",
            "dJ/du_o = (sigma(u_o^T v_c) - 1) v_c",
            "dJ/du_{w_i} = sigma(u_{w_i}^T v_c) v_c",
        ],
        [
            "Gradient for center word: pulled toward positive context, away from negatives.",
            "Gradient for positive context: pulled toward center word.",
            "Gradient for negative context: pushed away from center word.  "
            "Update: v <- v - alpha * dJ/dv  (subtract because we maximize J).",
        ], ACCENT_GREEN)

    add_table_slide(prs, "Noise Distribution: Effect of the 0.75 Exponent",
        ["Word", "Freq f(w)", "Unigram f(w)^1.0", "Smoothed f(w)^0.75", "Effect"],
        [
            ["the", "0.07", "0.07", "0.13", "Relatively reduced"],
            ["cat", "0.001", "0.001", "0.006", "Relatively boosted"],
            ["quasar", "0.00001", "0.00001", "0.0002", "Much more likely"],
        ], ACCENT_GREEN)

    # ── Section 5: Training Pipeline ──
    add_section_slide(prs, "5. SGNS Training Pipeline", ACCENT_GREEN)

    add_content_slide(prs, "SGNS Training Pipeline (6 Steps)", [
        "1. Preprocessing: Tokenize, subsample frequent words",
        "",
        "2. Initialize: Build vocabulary; create W (center) and C (context) matrices "
        "with small random values ~ Uniform(-0.5/d, 0.5/d)",
        "",
        "3. Select center word: Iterate through each word in corpus",
        "",
        "4. Generate training pairs:",
        "  - Positive: (center, context) pairs from the window",
        "  - Negative: Sample k random words from noise distribution P_n(w)",
        "",
        "5. Compute scores: Dot product through sigmoid; loss encourages high scores "
        "for positive pairs, low for negatives",
        "",
        "6. Sparse update via SGD: Only update embeddings for words in current sample "
        "(1 positive + k negatives) -- this is what makes SGNS efficient!",
    ], ACCENT_GREEN)

    # ── References ──
    add_section_slide(prs, "6. References", ACCENT_GREEN)
    add_content_slide(prs, "Key References", [
        "[1] Mikolov et al. (2013) - 'Distributed representations of words and phrases "
        "and their compositionality' (NIPS) -- Introduces negative sampling & subsampling",
        "",
        "[2] Mikolov et al. (2013) - 'Efficient estimation of word representations "
        "in vector space' (ICLR) -- Original Skip-gram & CBOW architectures",
        "",
        "[3] Goldberg & Levy (2014) - 'word2vec Explained' -- Clear mathematical derivation "
        "of SGNS objective and gradients",
        "",
        "[4] Levy & Goldberg (2014) - 'Neural word embedding as implicit matrix "
        "factorization' (NIPS) -- SGNS implicitly factorizes shifted PMI matrix",
        "",
        "[5] Gutmann & Hyvarinen (2010) - 'Noise-contrastive estimation' -- Theoretical "
        "foundation for negative sampling",
    ], ACCENT_GREEN)


# ══════════════════════════════════════════════════════════════════════════
# TOPIC 3: CBOW
# ══════════════════════════════════════════════════════════════════════════

def build_cbow(prs):
    add_title_slide(prs, "CBOW\nContinuous Bag of Words",
                    "Predicting Words from Context",
                    ACCENT_ORANGE)

    # ── Section 1: Intuition ──
    add_section_slide(prs, "1. Intuition & Core Concepts", ACCENT_ORANGE)

    add_key_concept_slide(prs, "What is CBOW?",
        "Given surrounding context words, predict the center word.",
        [
            "CBOW is a neural network model that learns word embeddings by predicting "
            "a center word from its surrounding context words.",
            'Like a fill-in-the-blank game: "The cat sat on the ___"',
            "Context words: 'The', 'cat', 'sat', 'on', 'the' --> Predict: 'mat'",
            "By learning to fill in blanks, CBOW discovers that words appearing in "
            "similar contexts have similar meanings.",
        ], ACCENT_ORANGE)

    add_content_slide(prs, 'Why "Bag of Words"?', [
        "CBOW averages context word vectors, ignoring their order.",
        "It treats context as a 'bag' (order-independent set of words).",
        "",
        "These contexts are treated identically:",
        '  - "quick brown ___ jumps over"',
        '  - "brown quick ___ over jumps"',
        "Both produce the same averaged context vector.",
        "",
        "This simplification makes CBOW fast and memory-efficient, "
        "though it loses positional information.",
    ], ACCENT_ORANGE)

    add_content_slide(prs, "CBOW Process Flow", [
        "1. INPUT: Context word embeddings (e.g., 'quick', 'brown', 'jumps')",
        "2. AVERAGE: Compute mean of all context vectors --> hidden vector h",
        '3. PREDICT: Use h to compute scores for every word in vocabulary',
        '4. SOFTMAX: Convert scores to probabilities (sum to 1)',
        '5. LOSS: Cross-entropy between predicted distribution and true word',
        '6. UPDATE: Adjust embeddings via backpropagation to improve predictions',
    ], ACCENT_ORANGE)

    # ── Section 2: Worked Example ──
    add_section_slide(prs, "2. Worked Example", ACCENT_ORANGE)

    add_content_slide(prs,
        'Processing "The quick brown fox jumps over the lazy dog"', [
        "Window size = 2.  Training examples:",
        "",
        "Position 0 (The):   Context [quick, brown]          --> Predict 'The'",
        "Position 1 (quick): Context [The, brown, fox]        --> Predict 'quick'",
        "Position 2 (brown): Context [The, quick, fox, jumps] --> Predict 'brown'",
        "Position 3 (fox):   Context [quick, brown, jumps, over] --> Predict 'fox'",
        "Position 4 (jumps): Context [brown, fox, over, the]  --> Predict 'jumps'",
        "  ...and so on for each position.",
    ], ACCENT_ORANGE)

    add_content_slide(prs,
        "Step-by-Step: Predict 'fox' from [quick, brown, jumps, over]", [
        "Context embeddings (dim=3):",
        "  quick = [0.2, 0.8, -0.1],  brown = [0.5, 0.3, 0.7]",
        "  jumps = [0.1, 0.9, 0.2],   over  = [0.4, 0.2, 0.5]",
        "",
        "Step 1 -- Average context vectors:",
        "  h[0] = (0.2+0.5+0.1+0.4)/4 = 0.30",
        "  h[1] = (0.8+0.3+0.9+0.2)/4 = 0.55",
        "  h[2] = (-0.1+0.7+0.2+0.5)/4 = 0.325",
        "  h = [0.30, 0.55, 0.325]",
        "",
        "Step 2 -- Compute dot product scores with output embeddings,",
        "apply softmax --> P(fox) = 0.35 (highest!)",
        "Step 3 -- Training adjusts embeddings to increase this probability.",
    ], ACCENT_ORANGE)

    # ── Section 3: Detailed Numerical Example ──
    add_section_slide(prs, "3. Full Numerical Example", ACCENT_ORANGE)

    add_content_slide(prs,
        'Predict "sat" from ["the", "cat", "on", "mat"]  (V=5, d=2)', [
        "Input Embeddings W:",
        "  the=[0.2,0.8], cat=[0.9,0.3], sat=[0.5,0.6], on=[0.1,0.7], mat=[0.8,0.4]",
        "",
        "Output Embeddings W':",
        "  the=[0.3,0.5], cat=[0.7,0.2], sat=[0.6,0.8], on=[0.2,0.6], mat=[0.5,0.3]",
        "",
        "Step 1 -- Average context: h = (v_the + v_cat + v_on + v_mat) / 4",
        "  h[0] = (0.2+0.9+0.1+0.8)/4 = 0.50",
        "  h[1] = (0.8+0.3+0.7+0.4)/4 = 0.55",
        "  h = [0.50, 0.55]",
    ], ACCENT_ORANGE)

    add_table_slide(prs, "Scores, Softmax, and Probabilities",
        ["Word", "Score (h . w')", "exp(score)", "Probability"],
        [
            ["the", "0.425", "1.530", "0.185"],
            ["cat", "0.460", "1.584", "0.192"],
            ["sat (target)", "0.740", "2.096", "0.254"],
            ["on", "0.430", "1.537", "0.186"],
            ["mat", "0.415", "1.514", "0.183"],
        ], ACCENT_ORANGE)

    add_content_slide(prs, "Loss Calculation and Interpretation", [
        "Model assigns highest probability (25.4%) to 'sat' -- the correct target!",
        "",
        "Cross-entropy loss:",
        "  L = -log P(sat) = -log(0.254) = 1.370",
        "",
        "Comparison:",
        "  Perfect prediction (P=1.0): Loss = 0",
        "  Our prediction (P=0.254):   Loss = 1.370",
        "  Random guess (P=0.2):       Loss = 1.609",
        "",
        "Model is better than random. Training continues to minimize this loss "
        "by adjusting both W and W' embeddings.",
    ], ACCENT_ORANGE)

    # ── Section 4: Mathematical Foundations ──
    add_section_slide(prs, "4. Mathematical Foundations", ACCENT_ORANGE)

    add_formula_slide(prs, "CBOW Objective Function",
        [
            "J(theta) = (1/T) SUM_{t=1}^{T} log P(w_t | w_{t-c}, ..., w_{t+c})",
        ],
        [
            "T = total words in corpus;  w_t = target (center) word;  "
            "c = context window size.  Goal: maximize probability of predicting "
            "the correct center word given its context.",
        ], ACCENT_ORANGE)

    add_formula_slide(prs, "Context Averaging & Softmax",
        [
            "h = (1/2c) SUM_{j, j!=0} v_{w_{t+j}}",
            "P(w_O | context) = exp(v'_{w_O}^T h) / SUM_{w=1}^{V} exp(v'_w^T h)",
        ],
        [
            "h = hidden layer vector = average of all context word embeddings (from W).  "
            "No learnable parameters in the averaging step -- keeps it simple.",
            "v'_{w_O} = output embedding for target word;  V = vocabulary size.  "
            "Denominator sums over all words (expensive -- use neg. sampling or hier. softmax).",
        ], ACCENT_ORANGE)

    add_formula_slide(prs, "Loss and Gradients",
        [
            "L = -log P(w_O | context) = -v'_{w_O}^T h + log SUM exp(v'_w^T h)",
            "dL/dv'_w = (P(w|context) - 1{w=w_O}) * h",
            "dL/dv_{w_i} = (1/2c) * dL/dh      (for each context word w_i)",
        ],
        [
            "Negative log-likelihood (cross-entropy) loss.",
            "For correct word: gradient pushes its output embedding closer to h.  "
            "For incorrect words: gradient pushes them away.",
            "Each context word's input embedding is updated by the same amount, "
            "scaled by 1/2c due to the averaging step.",
        ], ACCENT_ORANGE)

    add_table_slide(prs, "Two Sets of Embeddings",
        ["Matrix", "Symbol", "Shape", "Purpose"],
        [
            ["Input Embeddings", "W (or v)", "V x d", "Represent context words"],
            ["Output Embeddings", "W' (or v')", "V x d", "Compute prediction scores"],
        ], ACCENT_ORANGE)

    add_content_slide(prs, "Which Embeddings to Use After Training?", [
        "Option 1: Input embeddings W only (most common)",
        "Option 2: Output embeddings W' only (less common)",
        "Option 3: Average W and W' (sometimes improves quality)",
        "Option 4: Concatenate W and W' (doubles dimensionality)",
    ], ACCENT_ORANGE)

    # ── Section 5: CBOW vs Skip-gram ──
    add_section_slide(prs, "5. CBOW vs. Skip-gram", ACCENT_ORANGE)

    add_table_slide(prs, "Architecture Comparison",
        ["Aspect", "CBOW", "Skip-gram"],
        [
            ["Input", "Multiple context words", "Single center word"],
            ["Output", "Single center word", "Multiple context words"],
            ["Task", "P(center | context)", "P(context | center)"],
            ["Training Speed", "Faster (1 prediction/window)", "Slower (multiple/window)"],
            ["Frequent Words", "Better representations", "May underperform"],
            ["Rare Words", "May underperform", "Better representations"],
            ["Memory", "Lower (fewer gradients)", "Higher (more gradients)"],
        ], ACCENT_ORANGE)

    add_table_slide(prs, "Practical Recommendations",
        ["Scenario", "Recommendation", "Reason"],
        [
            ["Billions of words", "CBOW", "~4x faster training"],
            ["Small corpus, quality matters", "Skip-gram", "Better rare word handling"],
            ["Domain with many rare terms", "Skip-gram", "More training signal per word"],
            ["General-purpose embeddings", "Skip-gram + NS", "Best overall quality"],
            ["Online/streaming training", "CBOW", "Lower memory, faster updates"],
        ], ACCENT_ORANGE)

    # ── References ──
    add_section_slide(prs, "6. References", ACCENT_ORANGE)
    add_content_slide(prs, "Key References", [
        "[1] Mikolov et al. (2013) - 'Efficient estimation of word representations "
        "in vector space' (ICLR) -- Introduced both CBOW and Skip-gram",
        "",
        "[2] Mikolov et al. (2013) - 'Distributed representations of words and phrases "
        "and their compositionality' (NIPS) -- Negative sampling, subsampling",
        "",
        "[3] Rong (2014) - 'word2vec parameter learning explained' -- Detailed gradient "
        "derivations for CBOW and Skip-gram",
        "",
        "[4] Bengio et al. (2003) - 'A neural probabilistic language model' (JMLR) -- "
        "Foundation for distributed representations and neural language modeling",
        "",
        "[5] Collobert & Weston (2008) - 'A unified architecture for NLP' (ICML) -- "
        "Demonstrated neural embeddings across multiple NLP tasks",
    ], ACCENT_ORANGE)


# ══════════════════════════════════════════════════════════════════════════
# TOPIC 4: GloVe
# ══════════════════════════════════════════════════════════════════════════

def build_glove(prs):
    add_title_slide(prs, "GloVe\nGlobal Vectors for Word Representation",
                    "Learning Word Vectors from Co-occurrence Ratios\n"
                    "Pennington, Socher & Manning (2014)",
                    ACCENT_RED)

    # ── Section 1: Intuition ──
    add_section_slide(prs, "1. Intuition & Core Concepts", ACCENT_RED)

    add_key_concept_slide(prs, "What is GloVe?",
        "The ratio of co-occurrence probabilities encodes semantic "
        "relationships more effectively than raw counts alone.",
        [
            "GloVe learns word vectors by analyzing how often words co-occur "
            "across an entire text corpus.",
            "Core idea: It's not absolute co-occurrence counts that matter, "
            "but the ratios between them.",
            "GloVe bridges count-based methods (LSA) and prediction-based methods "
            "(Word2Vec), combining the best of both worlds.",
        ], ACCENT_RED)

    add_content_slide(prs, "The Social Network Analogy", [
        'Direct counting: "Alice spent time with Bob 50 times" -- but is that a lot?',
        "",
        "Ratios reveal meaning:",
        '  "Alice spends 10x more time with Bob than Carol,',
        '   but Charlie spends equal time with both"',
        "  -- now the relationships are clear!",
        "",
        "GloVe learns meaning from relative patterns of who 'hangs out' with whom "
        "(co-occurrence relationships).",
    ], ACCENT_RED)

    add_content_slide(prs, "The Co-occurrence Matrix", [
        'Core question: "Who appears near whom?"',
        "GloVe builds a matrix X where X[i,j] counts how often word j appears "
        "near word i across the entire corpus.",
        "",
        'Example: "The cat sat on the mat" (window=1):',
        '  "cat" co-occurs with "the" and "sat"',
        '  "sat" co-occurs with "cat" and "on"',
        "",
        "This global matrix captures corpus-wide statistics -- "
        "unlike Word2Vec which processes local windows one at a time.",
    ], ACCENT_RED)

    # ── Section 2: The Power of Ratios ──
    add_section_slide(prs, "2. The Power of Ratios", ACCENT_RED)

    add_table_slide(prs, 'Classic Example: "ice" vs. "steam"',
        ["Probe word k", "P(k|ice)", "P(k|steam)", "Ratio P(k|ice)/P(k|steam)",
         "Interpretation"],
        [
            ["solid", "High", "Low", ">> 1 (Large)", "Strongly ice-related"],
            ["gas", "Low", "High", "<< 1 (Small)", "Strongly steam-related"],
            ["water", "High", "High", "~ 1", "Related to both equally"],
            ["fashion", "Low", "Low", "~ 1", "Unrelated to both"],
        ], ACCENT_RED)

    add_content_slide(prs, "What the Ratios Encode", [
        "Ratio >> 1:  probe word k is more associated with 'ice' (like 'solid', 'cold')",
        "Ratio << 1:  probe word k is more associated with 'steam' (like 'gas', 'hot')",
        "Ratio ~ 1:   probe word k is equally related to both (like 'water')",
        "             or unrelated to both (like 'fashion')",
        "",
        "GloVe's key insight: Train vectors so that:",
        "  (w_ice - w_steam) . w_k  ~=  log( P(k|ice) / P(k|steam) )",
        "",
        "The difference vector captures exactly what distinguishes the two words!",
        "This is also why analogies like king - man + woman = queen work.",
    ], ACCENT_RED)

    add_two_column_slide(prs, "Best of Both Worlds",
        "Count-based Methods (LSA, HAL)",
        [
            "Use global corpus statistics",
            "Efficient computation",
            "Capture word co-occurrence",
            "",
            "Weakness: Miss nuanced relationships",
            "Limited on analogy tasks",
        ],
        "Prediction-based (Word2Vec)",
        [
            "Learn from local context windows",
            "Capture analogies well",
            "Flexible architecture",
            "",
            "Weakness: Don't use global stats directly",
            "Training sees each window once",
        ])

    add_key_concept_slide(prs, "GloVe: The Bridge",
        "w_i . w_j + b_i + b_j  =  log(X_ij)",
        [
            "GloVe learns word vectors so their dot product equals the log "
            "of their co-occurrence count.",
            "High co-occurrence = high dot product = similar vectors",
            "Low co-occurrence = low dot product = different vectors",
            "Log transform prevents very frequent pairs from dominating",
            "Uses global statistics (like LSA) but learns through optimization "
            "(like Word2Vec)",
        ], ACCENT_RED)

    # ── Section 3: Worked Example ──
    add_section_slide(prs, "3. Worked Example", ACCENT_RED)

    add_content_slide(prs, 'Tiny Corpus (window=1)', [
        'Sentences:  "ice water is cold",  "steam water is hot",',
        '            "ice is solid",  "steam is gas"',
        "",
        "Co-occurrence matrix (selected entries):",
        "  X[ice, water]=1, X[ice, solid]=1, X[steam, hot]=1, X[steam, gas]=1",
        "  X[water, ice]=1, X[water, steam]=1  (water connects both)",
        "  X[is, ...] = high counts (appears in all sentences)",
    ], ACCENT_RED)

    add_table_slide(prs, "Co-occurrence Probabilities and Ratios",
        ["Probe k", "P(k|ice)", "P(k|steam)", "Ratio", "Meaning"],
        [
            ["solid", "0.333", "0.000", "Large (inf)", "Strongly ice-related"],
            ["gas", "0.000", "0.333", "Small (~0)", "Strongly steam-related"],
            ["water", "0.333", "0.333", "1.0", "Related to both equally"],
            ["is", "0.333", "0.333", "1.0", "Common word, neutral"],
        ], ACCENT_RED)

    add_content_slide(prs, "What GloVe Learns from These Ratios", [
        "Vector relationship goal:",
        "  (w_ice - w_steam) . w_k  ~=  log( P(k|ice) / P(k|steam) )",
        "",
        "The difference vector (w_ice - w_steam) should have:",
        '  Positive dot product with "solid" (ice-related)',
        '  Negative dot product with "gas" (steam-related)',
        '  Near zero dot product with "water" (both-related)',
        "",
        "This is also why king - man + woman = queen works!",
        "Co-occurrence ratios encode gender relationships consistently.",
    ], ACCENT_RED)

    # ── Section 4: Mathematical Foundations ──
    add_section_slide(prs, "4. Mathematical Foundations", ACCENT_RED)

    add_formula_slide(prs, "Co-occurrence Probability",
        [
            "P_ij = P(j|i) = X_ij / X_i     where X_i = SUM_k X_ik",
        ],
        [
            "X = word-word co-occurrence matrix;  X_ij = count of word j in "
            "context of word i;  X_i = total co-occurrences for word i.",
        ], ACCENT_RED)

    add_formula_slide(prs, "The GloVe Objective Function",
        [
            "J = SUM_{i,j=1}^{V}  f(X_ij) * (w_i^T w~_j + b_i + b~_j - log X_ij)^2",
        ],
        [
            "Weighted least squares: minimize squared error between dot product "
            "and log co-occurrence.  f(X_ij) is a weighting function.  "
            "Only sums over non-zero entries (sparse matrix).  "
            "V = vocabulary size.",
        ], ACCENT_RED)

    add_formula_slide(prs, "The Weighting Function f(x)",
        [
            "f(x) = (x / x_max)^alpha   if x < x_max,     1   otherwise",
        ],
        [
            "Typical values: x_max = 100, alpha = 0.75.  "
            "Rare co-occurrences get lower weight (might be noise).  "
            "Frequent co-occurrences capped at weight 1 (don't over-emphasize).  "
            "alpha = 0.75 provides sublinear scaling.",
        ], ACCENT_RED)

    add_table_slide(prs, "GloVe Model Components",
        ["Component", "Symbol", "Dimension", "Purpose"],
        [
            ["Word vectors", "W", "V x d", "Main embeddings"],
            ["Context vectors", "W~", "V x d", "Separate context embeddings"],
            ["Word biases", "b", "V x 1", "Capture word frequency effects"],
            ["Context biases", "b~", "V x 1", "Capture context frequency effects"],
            ["Co-occurrence matrix", "X", "V x V", "Input statistics from corpus"],
        ], ACCENT_RED)

    # ── Section 5: One Training Iteration ──
    add_section_slide(prs, "5. One Training Iteration (Numerical)", ACCENT_RED)

    add_content_slide(prs, "Setup: Training on pair (ice, cold), X[ice,cold]=1", [
        "2D vectors for simplicity:",
        "  w_ice = [0.5, -0.3],  w_ctx_cold = [0.3, -0.2]",
        "  b_ice = 0.1,  b_ctx_cold = 0.02",
        "",
        "Step 1 -- Forward pass:",
        "  Dot product: [0.5,-0.3].[0.3,-0.2] = 0.15+0.06 = 0.21",
        "  Prediction: 0.21 + 0.1 + 0.02 = 0.33",
        "",
        "Step 2 -- Target and error:",
        "  Target: log(X[ice,cold]) = log(1) = 0.0",
        "  Error (diff): 0.33 - 0.0 = 0.33",
    ], ACCENT_RED)

    add_content_slide(prs, "Weighting, Gradients, and Update", [
        "Step 3 -- Weighted loss:",
        "  f(1) = (1/100)^0.75 = 0.0316",
        "  L = 0.0316 * (0.33)^2 = 0.00344",
        "",
        "Step 4 -- Gradients (grad_common = 2 * 0.0316 * 0.33 = 0.0209):",
        "  grad_w_ice = 0.0209 * [0.3, -0.2] = [0.00627, -0.00418]",
        "  grad_w_ctx_cold = 0.0209 * [0.5, -0.3] = [0.01045, -0.00627]",
        "  grad_biases = 0.0209",
        "",
        "Step 5 -- SGD update (lr=0.05):",
        "  w_ice: [0.5, -0.3] - 0.05*[0.00627, -0.00418] = [0.4997, -0.2998]",
        "  b_ice: 0.1 - 0.05*0.0209 = 0.0990",
        "",
        "Prediction was too high (0.33 vs target 0.0), so vectors nudged apart slightly.",
    ], ACCENT_RED)

    # ── Section 6: Training Pipeline ──
    add_section_slide(prs, "6. Training Pipeline", ACCENT_RED)

    add_content_slide(prs, "GloVe Training Pipeline (5 Steps)", [
        "1. Preprocessing: Tokenize, lowercase, build vocabulary",
        "",
        "2. Build co-occurrence matrix: Scan corpus with context window "
        "(typically 10-15 words); weight by distance (closer = higher count)",
        "",
        "3. Initialize: Random word vectors W, context vectors W~, "
        "biases b, b~;  typical d = 50-300",
        "",
        "4. Training loop (until convergence):",
        "  - Forward: compute dot product + biases",
        "  - Loss: weighted squared error vs log(X_ij)",
        "  - Backward: compute gradients",
        "  - Update: AdaGrad optimizer (adapts learning rate per parameter)",
        "",
        "5. Final embeddings: Use W alone, or W + W~ (often slightly better)",
    ], ACCENT_RED)

    # ── Section 7: GloVe vs Word2Vec ──
    add_section_slide(prs, "7. GloVe vs. Word2Vec", ACCENT_RED)

    add_table_slide(prs, "GloVe vs. Word2Vec Comparison",
        ["Aspect", "GloVe", "Word2Vec (Skip-gram)"],
        [
            ["Training Approach", "Factorizes co-occurrence matrix",
             "Predicts context words"],
            ["Statistics Used", "Global corpus statistics", "Local context windows"],
            ["Training Data", "Pre-computed co-occurrence counts",
             "Raw text (online learning)"],
            ["Memory", "Stores co-occurrence matrix", "Streams through corpus"],
            ["Parallelization", "Easy (matrix factorization)",
             "Requires special handling"],
            ["Analogy Tasks", "Excellent", "Excellent"],
            ["Rare Words", "Explicit weighting function",
             "Negative sampling handles this"],
        ], ACCENT_RED)

    add_table_slide(prs, "Embedding Methods Comparison",
        ["Method", "What It Learns", "Best For"],
        [
            ["TF-IDF", "Word importance in documents",
             "Information retrieval, keyword extraction"],
            ["Word2Vec", "Predictive word relationships",
             "General semantic similarity"],
            ["GloVe", "Co-occurrence ratios",
             "Analogy tasks, semantic relationships"],
            ["FastText", "Subword + word embeddings",
             "Morphologically rich languages, rare words"],
            ["BERT", "Contextualized embeddings",
             "Context-dependent meaning"],
        ], ACCENT_RED)

    # ── References ──
    add_section_slide(prs, "8. References", ACCENT_RED)
    add_content_slide(prs, "Key References", [
        "[1] Pennington, Socher & Manning (2014) - 'GloVe: Global Vectors for "
        "Word Representation' (EMNLP) -- The original GloVe paper",
        "",
        "[2] Levy & Goldberg (2014) - 'Neural word embedding as implicit matrix "
        "factorization' (NIPS) -- Connects SGNS to shifted PMI matrix factorization",
        "",
        "[3] Levy, Goldberg & Dagan (2015) - 'Improving distributional similarity "
        "with lessons learned from word embeddings' (TACL) -- Systematic comparison",
        "",
        "[4] Baroni, Dinu & Kruszewski (2014) - \"Don't count, predict!\" (ACL) -- "
        "Count-based vs. prediction-based comparison motivating GloVe",
        "",
        "[5] Turney & Pantel (2010) - 'From frequency to meaning: Vector space models "
        "of semantics' (JAIR) -- Comprehensive survey of vector space models",
    ], ACCENT_RED)


# ══════════════════════════════════════════════════════════════════════════
# MAIN
# ══════════════════════════════════════════════════════════════════════════

def main() -> None:
    """Generate 4 PPTX files, one per topic."""
    output_dir = "../resources"

    topics = [
        ("word2vec", build_word2vec),
        ("sgns", build_sgns),
        ("cbow", build_cbow),
        ("glove", build_glove),
    ]

    for name, builder in topics:
        prs = _new_prs()
        builder(prs)
        path = f"{output_dir}/{name}_presentation.pptx"
        prs.save(path)
        print(f"Saved: {path}  ({len(prs.slides)} slides)")


if __name__ == "__main__":
    main()
