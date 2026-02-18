"""Generate a sample PPTX quiz deck for review before producing all 4 topics."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE


def hex_to_rgb(hex_str: str) -> RGBColor:
    """Convert a hex color string to RGBColor."""
    h = hex_str.lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


# Color palette
BG_DARK = hex_to_rgb("1B2A4A")
BG_CARD = hex_to_rgb("223556")
ACCENT_BLUE = hex_to_rgb("4A90D9")
ACCENT_GREEN = hex_to_rgb("27AE60")
TEXT_WHITE = hex_to_rgb("FFFFFF")
TEXT_LIGHT = hex_to_rgb("D0D8E8")
TEXT_MUTED = hex_to_rgb("8899AA")
CORRECT_BG = hex_to_rgb("1E3A2F")
CORRECT_BORDER = hex_to_rgb("27AE60")
WRONG_BG = hex_to_rgb("3A1E1E")
WRONG_BORDER = hex_to_rgb("C0392B")
OPTION_BG = hex_to_rgb("2A3F66")


def set_slide_bg(slide, color: RGBColor) -> None:
    """Set the background color of a slide."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rounded_rect(slide, left, top, width, height, fill_color, line_color=None):
    """Add a rounded rectangle shape."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.fill.solid()
        shape.line.fill.fore_color.rgb = line_color
        shape.line.width = Pt(2)
    else:
        shape.line.fill.background()
    return shape


def set_text(shape, text: str, font_size: int = 14, color: RGBColor = TEXT_WHITE,
             bold: bool = False, alignment=PP_ALIGN.LEFT) -> None:
    """Set text on a shape's text frame."""
    tf = shape.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.alignment = alignment


def create_title_slide(prs: Presentation, topic_num: int, title: str,
                       subtitle: str) -> None:
    """Create a title slide for a topic."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    set_slide_bg(slide, BG_DARK)

    # Topic number badge
    badge = add_rounded_rect(
        slide, Inches(3.8), Inches(1.0), Inches(2.4), Inches(0.6), ACCENT_BLUE
    )
    set_text(badge, f"TOPIC {topic_num}", 16, TEXT_WHITE, True, PP_ALIGN.CENTER)

    # Title
    title_box = slide.shapes.add_textbox(
        Inches(1.0), Inches(2.0), Inches(8.0), Inches(1.2)
    )
    set_text(title_box, title, 32, TEXT_WHITE, True, PP_ALIGN.CENTER)

    # Subtitle
    sub_box = slide.shapes.add_textbox(
        Inches(1.5), Inches(3.3), Inches(7.0), Inches(0.8)
    )
    set_text(sub_box, subtitle, 16, TEXT_MUTED, False, PP_ALIGN.CENTER)

    # Decorative line
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(3.5), Inches(4.3), Inches(3.0), Pt(3)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = ACCENT_BLUE
    line.line.fill.background()


def create_question_slide(prs: Presentation, q_num: int, question: str,
                          options: list, correct_indices: list,
                          show_answers: bool = False) -> None:
    """Create a question slide.

    Args:
        prs: Presentation object.
        q_num: Question number.
        question: Question text.
        options: List of option strings (A, B, C, D).
        correct_indices: List of 0-based indices of correct answers.
        show_answers: If True, highlight correct/incorrect answers.
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    set_slide_bg(slide, BG_DARK)

    # Question number badge
    badge = add_rounded_rect(
        slide, Inches(0.4), Inches(0.3), Inches(1.6), Inches(0.45), ACCENT_BLUE
    )
    set_text(badge, f"Question {q_num}", 13, TEXT_WHITE, True, PP_ALIGN.CENTER)

    # "Select all that apply" tag
    tag = add_rounded_rect(
        slide, Inches(7.5), Inches(0.3), Inches(2.1), Inches(0.45), OPTION_BG
    )
    set_text(tag, "Select all that apply", 11, TEXT_MUTED, False, PP_ALIGN.CENTER)

    # Question text
    q_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(1.0), Inches(9.0), Inches(1.2)
    )
    tf = q_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = question
    p.font.size = Pt(20)
    p.font.color.rgb = TEXT_WHITE
    p.font.bold = True

    # Options
    labels = ["A", "B", "C", "D"]
    top_start = 2.5
    option_height = 0.7
    option_gap = 0.15

    for i, option_text in enumerate(options):
        top = Inches(top_start + i * (option_height + option_gap))
        is_correct = i in correct_indices

        if show_answers:
            bg = CORRECT_BG if is_correct else WRONG_BG
            border = CORRECT_BORDER if is_correct else WRONG_BORDER
        else:
            bg = OPTION_BG
            border = None

        # Option card
        card = add_rounded_rect(
            slide, Inches(0.8), top, Inches(8.4), Inches(option_height),
            bg, border
        )

        # Label circle
        circle = slide.shapes.add_shape(
            MSO_SHAPE.OVAL, Inches(1.0), top + Inches(0.12),
            Inches(0.45), Inches(0.45)
        )
        circle.fill.solid()
        if show_answers and is_correct:
            circle.fill.fore_color.rgb = ACCENT_GREEN
        else:
            circle.fill.fore_color.rgb = ACCENT_BLUE
        circle.line.fill.background()
        set_text(circle, labels[i], 14, TEXT_WHITE, True, PP_ALIGN.CENTER)

        # Option text
        opt_box = slide.shapes.add_textbox(
            Inches(1.7), top + Inches(0.1), Inches(7.2), Inches(0.5)
        )
        tf = opt_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = option_text
        p.font.size = Pt(15)
        p.font.color.rgb = TEXT_WHITE if not show_answers else (
            CORRECT_BORDER if is_correct else hex_to_rgb("E07070")
        )

        # Checkmark or X for answer slides
        if show_answers:
            marker = slide.shapes.add_textbox(
                Inches(8.6), top + Inches(0.1), Inches(0.5), Inches(0.5)
            )
            p = marker.text_frame.paragraphs[0]
            p.text = "\u2713" if is_correct else "\u2717"
            p.font.size = Pt(20)
            p.font.bold = True
            p.font.color.rgb = CORRECT_BORDER if is_correct else hex_to_rgb("C0392B")


def main() -> None:
    """Generate a sample PPTX with a few slides for review."""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)  # 16:9

    # Title slide
    create_title_slide(
        prs, 1,
        "NLU/NLG, NLP Applications\n& Levels of Analysis",
        "NLP Knowledge Base \u2014 Quiz Review"
    )

    # Sample Question (without answers)
    create_question_slide(
        prs, 1,
        "Which of the following are Natural Language Understanding (NLU) tasks?",
        [
            "Generating a weather report from structured data",
            "Classifying an email as spam or not spam",
            "Producing a summary paragraph from bullet points",
            "Determining the sentiment of a movie review",
        ],
        correct_indices=[1, 3],
        show_answers=False,
    )

    # Same question with answers revealed
    create_question_slide(
        prs, 1,
        "Which of the following are Natural Language Understanding (NLU) tasks?",
        [
            "Generating a weather report from structured data",
            "Classifying an email as spam or not spam",
            "Producing a summary paragraph from bullet points",
            "Determining the sentiment of a movie review",
        ],
        correct_indices=[1, 3],
        show_answers=True,
    )

    # Another sample question (without answers)
    create_question_slide(
        prs, 2,
        "Which types of ambiguity can make NLP difficult?",
        [
            "Lexical ambiguity \u2014 a word has multiple meanings",
            "Syntactic ambiguity \u2014 a sentence has multiple valid parse trees",
            "Alphabetical ambiguity \u2014 letters can appear in multiple fonts",
            "Referential ambiguity \u2014 a pronoun could refer to multiple entities",
        ],
        correct_indices=[0, 1, 3],
        show_answers=False,
    )

    # Same with answers
    create_question_slide(
        prs, 2,
        "Which types of ambiguity can make NLP difficult?",
        [
            "Lexical ambiguity \u2014 a word has multiple meanings",
            "Syntactic ambiguity \u2014 a sentence has multiple valid parse trees",
            "Alphabetical ambiguity \u2014 letters can appear in multiple fonts",
            "Referential ambiguity \u2014 a pronoun could refer to multiple entities",
        ],
        correct_indices=[0, 1, 3],
        show_answers=True,
    )

    output_path = "../resources/nlp_quiz_sample.pptx"
    prs.save(output_path)
    print(f"Sample saved to {output_path}")


if __name__ == "__main__":
    main()
