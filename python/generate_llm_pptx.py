"""
Generate LLMs & Prompt Engineering PPTX presentation.
Run: python python/generate_llm_pptx.py
Output: html/llm_prompt_engineering.pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import os

# ── Colour palette ────────────────────────────────────────────────────────────
DARK_BG      = RGBColor(0x0F, 0x17, 0x2A)   # deep navy
ACCENT_BLUE  = RGBColor(0x38, 0xBD, 0xF8)   # sky blue
ACCENT_GREEN = RGBColor(0x34, 0xD3, 0x99)   # emerald
ACCENT_PURP  = RGBColor(0xA7, 0x8B, 0xFA)   # violet
ACCENT_AMBER = RGBColor(0xFB, 0xBF, 0x24)   # amber
ACCENT_ROSE  = RGBColor(0xF8, 0x71, 0x71)   # rose
WHITE        = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY   = RGBColor(0xCB, 0xD5, 0xE1)
MID_GRAY     = RGBColor(0x64, 0x74, 0x8B)
CARD_BG      = RGBColor(0x1E, 0x29, 0x3B)   # slightly lighter navy
SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)

# ── Helpers ───────────────────────────────────────────────────────────────────

def new_prs():
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs


def blank_slide(prs):
    layout = prs.slide_layouts[6]  # completely blank
    return prs.slides.add_slide(layout)


def fill_bg(slide, color=DARK_BG):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, x, y, w, h, fill_color, alpha=None):
    shape = slide.shapes.add_shape(1, x, y, w, h)  # MSO_SHAPE_TYPE.RECTANGLE
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape


def add_text(slide, text, x, y, w, h, font_size=18, bold=False, color=WHITE,
             align=PP_ALIGN.LEFT, wrap=True, italic=False):
    txBox = slide.shapes.add_textbox(x, y, w, h)
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txBox


def add_para(tf, text, font_size=16, bold=False, color=WHITE,
             align=PP_ALIGN.LEFT, space_before=Pt(4), italic=False, bullet=False):
    p = tf.add_paragraph()
    p.alignment = align
    if space_before:
        p.space_before = space_before
    if bullet:
        p.level = 1
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return p


def slide_header(slide, title, subtitle=None, accent=ACCENT_BLUE):
    """Top accent bar + title + optional subtitle."""
    # accent bar
    add_rect(slide, 0, 0, SLIDE_W, Inches(0.08), accent)
    # title
    add_text(slide, title,
             Inches(0.5), Inches(0.18), Inches(12.3), Inches(0.7),
             font_size=32, bold=True, color=WHITE)
    if subtitle:
        add_text(slide, subtitle,
                 Inches(0.5), Inches(0.85), Inches(12.3), Inches(0.45),
                 font_size=17, color=LIGHT_GRAY)


def section_divider(prs, title, subtitle="", color=ACCENT_BLUE):
    """Full-bleed section-break slide."""
    slide = blank_slide(prs)
    fill_bg(slide, CARD_BG)
    add_rect(slide, 0, Inches(2.8), SLIDE_W, Inches(1.9), color)
    add_text(slide, title,
             Inches(0.8), Inches(2.95), Inches(11.7), Inches(1.2),
             font_size=42, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    if subtitle:
        add_text(slide, subtitle,
                 Inches(0.8), Inches(4.15), Inches(11.7), Inches(0.6),
                 font_size=20, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)
    return slide


def bullet_card(slide, x, y, w, h, title, bullets, title_color=ACCENT_BLUE,
                bg=CARD_BG, font_size=14):
    add_rect(slide, x, y, w, h, bg)
    # thin left border
    add_rect(slide, x, y, Inches(0.06), h, title_color)
    add_text(slide, title, x+Inches(0.12), y+Inches(0.08),
             w-Inches(0.2), Inches(0.38),
             font_size=15, bold=True, color=title_color)
    txBox = slide.shapes.add_textbox(x+Inches(0.12), y+Inches(0.5),
                                     w-Inches(0.24), h-Inches(0.6))
    tf = txBox.text_frame
    tf.word_wrap = True
    first = True
    for b in bullets:
        if first:
            p = tf.paragraphs[0]; first = False
        else:
            p = tf.add_paragraph()
        p.space_before = Pt(3)
        run = p.add_run()
        run.text = ("• " if not b.startswith("→") else "") + b
        run.font.size = Pt(font_size)
        run.font.color.rgb = LIGHT_GRAY


def tag_badge(slide, x, y, text, color=ACCENT_BLUE):
    w = Inches(1.6); h = Inches(0.32)
    add_rect(slide, x, y, w, h, color)
    add_text(slide, text, x, y, w, h,
             font_size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)


# ── SLIDE BUILDERS ────────────────────────────────────────────────────────────

def s_title(prs):
    slide = blank_slide(prs)
    fill_bg(slide)
    # gradient-like left stripe
    add_rect(slide, 0, 0, Inches(0.5), SLIDE_H, ACCENT_BLUE)
    add_rect(slide, 0, 0, Inches(0.5), Inches(2.5), ACCENT_PURP)
    add_rect(slide, 0, Inches(5.0), Inches(0.5), Inches(2.5), ACCENT_GREEN)

    add_text(slide, "Large Language Models",
             Inches(0.9), Inches(1.1), Inches(11.5), Inches(1.0),
             font_size=46, bold=True, color=WHITE)
    add_text(slide, "& Prompt Engineering",
             Inches(0.9), Inches(2.0), Inches(11.5), Inches(0.9),
             font_size=40, bold=True, color=ACCENT_BLUE)
    add_text(slide,
             "From transformer basics to RAG, tool use, hallucination mitigation,\n"
             "and production evaluation frameworks",
             Inches(0.9), Inches(3.1), Inches(10.0), Inches(1.0),
             font_size=18, color=LIGHT_GRAY)
    add_text(slide, "NLP Knowledge Base  ·  Educational Series",
             Inches(0.9), Inches(6.7), Inches(8.0), Inches(0.4),
             font_size=13, color=MID_GRAY)


def s_agenda(prs):
    slide = blank_slide(prs)
    fill_bg(slide)
    slide_header(slide, "Agenda", accent=ACCENT_PURP)
    topics = [
        ("①", "What is an LLM?",              "Architecture, scale, tokenisation"),
        ("②", "Training Pipeline",             "Pre-training → Fine-tuning → RLHF"),
        ("③", "Prompt Engineering 101",        "Zero-shot, Few-shot, anatomy"),
        ("④", "Chain-of-Thought & Variants",   "CoT, Least-to-Most, Skeleton, GoT"),
        ("⑤", "Generated Knowledge",           "Two-step knowledge injection"),
        ("⑥", "RAG",                           "Retrieval-Augmented Generation"),
        ("⑦", "Function Calling / Tool Use",   "JSON schema, dispatch loop"),
        ("⑧", "Hallucination & Safety",        "Types, mitigations, injection"),
        ("⑨", "Evaluation & Iteration",        "Metrics, A/B testing, fine-tune vs few-shot"),
        ("⑩", "Key Takeaways",                 "Production checklist"),
    ]
    cols = 2
    per_col = (len(topics) + 1) // cols
    for i, (num, title, desc) in enumerate(topics):
        col = i // per_col
        row = i % per_col
        x = Inches(0.5 + col * 6.4)
        y = Inches(1.4 + row * 0.58)
        add_text(slide, num,  x, y, Inches(0.5), Inches(0.5),
                 font_size=18, bold=True, color=ACCENT_BLUE)
        add_text(slide, title, x+Inches(0.5), y, Inches(5.0), Inches(0.32),
                 font_size=16, bold=True, color=WHITE)
        add_text(slide, desc,  x+Inches(0.5), y+Inches(0.3), Inches(5.5), Inches(0.26),
                 font_size=12, color=MID_GRAY)


def s_what_is_llm(prs):
    slide = blank_slide(prs)
    fill_bg(slide)
    slide_header(slide, "What is a Large Language Model?", accent=ACCENT_BLUE)

    # Left column: definition + analogy
    bullet_card(slide, Inches(0.4), Inches(1.25), Inches(5.8), Inches(2.3),
                "Core Idea",
                ["Trained on hundreds of billions of tokens of text",
                 "Learns statistical patterns → predicts next token",
                 "Emergent capabilities arise at scale (reasoning, analogies, code)",
                 "Not a database — probabilistic, not deterministic"],
                title_color=ACCENT_BLUE)

    bullet_card(slide, Inches(0.4), Inches(3.65), Inches(5.8), Inches(2.2),
                "The Well-Read Student Analogy",
                ["Reads everything (Wikipedia, books, code, papers)",
                 "Infers rules without being explicitly taught grammar",
                 "Can answer new questions by analogy to training",
                 "May 'confabulate' when out of distribution"],
                title_color=ACCENT_GREEN)

    # Right column: model scale table
    headers = ["Model", "Params", "Year", "Context"]
    rows = [
        ["GPT-3",     "175 B",  "2020", "4 K"],
        ["PaLM",      "540 B",  "2022", "8 K"],
        ["GPT-4",     "~1.7 T", "2023", "128 K"],
        ["Llama 3",   "70 B",   "2024", "128 K"],
        ["Claude 3.5","~100 B", "2024", "200 K"],
        ["Gemini 1.5","~1 T",   "2024", "1 M"],
    ]
    x0 = Inches(6.6); y0 = Inches(1.25)
    col_ws = [Inches(2.2), Inches(1.3), Inches(1.0), Inches(1.5)]
    row_h = Inches(0.38)
    # header row
    for ci, (hdr, cw) in enumerate(zip(headers, col_ws)):
        cx = x0 + sum(col_ws[:ci])
        add_rect(slide, cx, y0, cw, row_h, ACCENT_BLUE)
        add_text(slide, hdr, cx+Inches(0.05), y0+Inches(0.05),
                 cw-Inches(0.1), row_h-Inches(0.06),
                 font_size=12, bold=True, color=WHITE)
    for ri, row in enumerate(rows):
        bg = CARD_BG if ri % 2 == 0 else RGBColor(0x16, 0x20, 0x33)
        for ci, (cell, cw) in enumerate(zip(row, col_ws)):
            cx = x0 + sum(col_ws[:ci])
            cy = y0 + (ri+1)*row_h
            add_rect(slide, cx, cy, cw, row_h, bg)
            add_text(slide, cell, cx+Inches(0.05), cy+Inches(0.05),
                     cw-Inches(0.1), row_h-Inches(0.06),
                     font_size=11, color=LIGHT_GRAY)


def s_transformer(prs):
    slide = blank_slide(prs)
    fill_bg(slide)
    slide_header(slide, "Transformer Architecture & Self-Attention", accent=ACCENT_PURP)

    bullet_card(slide, Inches(0.4), Inches(1.25), Inches(5.8), Inches(5.5),
                "Key Components",
                ["Tokeniser: BPE splits 'unhappiness' → 'un', 'happi', 'ness'",
                 "Token Embeddings: each token → dense vector (e.g. 4096-dim)",
                 "Positional Encoding: injects token order information",
                 "Multi-Head Self-Attention: every token attends to every other",
                 "  Q·Kᵀ / √d_k  →  softmax  →  weighted sum of V",
                 "Feed-Forward Network: per-token MLP (2-layer, ReLU/GELU)",
                 "Layer Norm + Residual connections: stabilise training",
                 "N stacked layers (GPT-3: 96 layers, 96 heads, d=12288)",
                 "Causal masking: token i can only attend to tokens ≤ i"],
                title_color=ACCENT_PURP, font_size=13)

    # Self-attention visual (ASCII-style boxes)
    boxes = [
        (Inches(7.0), Inches(1.4),  "Input Tokens",     ACCENT_BLUE),
        (Inches(7.0), Inches(2.3),  "Q  K  V Projections", ACCENT_PURP),
        (Inches(7.0), Inches(3.2),  "Scaled Dot-Product\nAttention", ACCENT_GREEN),
        (Inches(7.0), Inches(4.3),  "Feed-Forward\nNetwork", ACCENT_AMBER),
        (Inches(7.0), Inches(5.3),  "Output Logits → Softmax", ACCENT_ROSE),
    ]
    for bx, by, label, color in boxes:
        add_rect(slide, bx, by, Inches(5.7), Inches(0.72), CARD_BG)
        add_rect(slide, bx, by, Inches(0.08), Inches(0.72), color)
        add_text(slide, label, bx+Inches(0.15), by+Inches(0.12),
                 Inches(5.4), Inches(0.58), font_size=13, color=WHITE)
    # arrows between boxes
    for i in range(len(boxes)-1):
        _, by, _, _ = boxes[i]
        add_text(slide, "↓", Inches(9.8), by+Inches(0.68),
                 Inches(0.4), Inches(0.4), font_size=16, color=MID_GRAY,
                 align=PP_ALIGN.CENTER)


def s_training_pipeline(prs):
    slide = blank_slide(prs)
    fill_bg(slide)
    slide_header(slide, "LLM Training Pipeline", accent=ACCENT_GREEN)

    stages = [
        ("1  Pre-training",  ACCENT_BLUE,
         ["Objective: predict next token (causal LM)",
          "Data: Common Crawl, Books, GitHub, Wikipedia (trillions of tokens)",
          "Compute: thousands of A100s for weeks/months",
          "Result: strong text completion, broad world knowledge"]),
        ("2  Fine-tuning (SFT)", ACCENT_GREEN,
         ["Objective: learn task format from (prompt, response) pairs",
          "Data: 10K–1M human-written examples",
          "Compute: hours to days on a few GPUs",
          "Result: follows instructions, stays on topic"]),
        ("3  RLHF / RLAIF", ACCENT_PURP,
         ["Objective: align with human preferences",
          "Step A: train a reward model on human comparisons",
          "Step B: PPO/DPO — update LLM to maximise reward",
          "Result: helpful, harmless, honest outputs"]),
    ]
    for i, (title, color, bullets) in enumerate(stages):
        x = Inches(0.4 + i * 4.3)
        bullet_card(slide, x, Inches(1.25), Inches(4.1), Inches(5.5),
                    title, bullets, title_color=color, font_size=13)
    # arrows between stages
    for i in range(2):
        ax = Inches(4.35 + i * 4.3)
        add_text(slide, "→", ax, Inches(3.6), Inches(0.5), Inches(0.5),
                 font_size=28, bold=True, color=ACCENT_AMBER, align=PP_ALIGN.CENTER)


def s_prompt_anatomy(prs):
    slide = blank_slide(prs)
    fill_bg(slide)
    slide_header(slide, "Prompt Engineering — Anatomy of a Good Prompt", accent=ACCENT_AMBER)

    parts = [
        ("System Prompt",    ACCENT_BLUE,   Inches(1.25),
         "You are an expert data analyst. Answer concisely and cite sources.",
         "Sets role, tone, constraints, output format"),
        ("Context / Examples", ACCENT_GREEN, Inches(2.55),
         "Here is the dataset summary:\n  Revenue Q1: $1.2M, Q2: $1.8M\nExample: Q: 'What grew?' A: 'Revenue grew 50% QoQ'",
         "Grounds the model; few-shot examples here"),
        ("User Query",       ACCENT_PURP,   Inches(4.1),
         "What was the growth rate between Q1 and Q2?",
         "Clear, specific, single task"),
        ("Output Format",    ACCENT_AMBER,  Inches(5.15),
         'Respond in JSON: {"answer": "...", "confidence": "high/medium/low"}',
         "Structured output → easy parsing"),
    ]
    for title, color, y, example, note in parts:
        bullet_card(slide, Inches(0.4), y, Inches(5.5), Inches(1.05),
                    title, [example], title_color=color, font_size=12)
        add_text(slide, f"↳ {note}", Inches(6.1), y+Inches(0.3),
                 Inches(6.8), Inches(0.5), font_size=12, color=MID_GRAY, italic=True)

    add_text(slide, "Four principles: Be Specific  ·  Provide Context  ·  Structure Output  ·  Use Delimiters",
             Inches(0.4), Inches(6.6), Inches(12.5), Inches(0.45),
             font_size=14, bold=True, color=ACCENT_AMBER, align=PP_ALIGN.CENTER)


def s_nshot(prs):
    slide = blank_slide(prs)
    fill_bg(slide)
    slide_header(slide, "N-Shot Prompting", accent=ACCENT_BLUE)

    cards = [
        ("Zero-Shot", ACCENT_BLUE,
         ["No examples — relies entirely on pre-training",
          "Fast, cheap, good baseline",
          "Example: 'Classify as positive/negative: Great movie!'",
          "Best for: simple classification, summarisation",
          "Risk: ambiguous format, model may not follow intent"]),
        ("One-Shot", ACCENT_GREEN,
         ["Single input/output example",
          "Clarifies expected output format",
          "Example: 'Input: Bad film → Label: negative\\nInput: Great movie!'",
          "Best for: format-sensitive tasks",
          "Risk: single example may bias the model"]),
        ("Few-Shot (3–8)", ACCENT_PURP,
         ["Multiple diverse examples",
          "Reduces format ambiguity dramatically",
          "Balance labels: equal pos/neg examples",
          "Order matters: put hardest cases last",
          "Best for: classification, extraction, translation"]),
    ]
    for i, (title, color, bullets) in enumerate(cards):
        bullet_card(slide, Inches(0.4 + i * 4.3), Inches(1.3),
                    Inches(4.1), Inches(5.4),
                    title, bullets, title_color=color, font_size=13)


def s_cot(prs):
    slide = blank_slide(prs)
    fill_bg(slide)
    slide_header(slide, "Chain-of-Thought Prompting", accent=ACCENT_GREEN)

    bullet_card(slide, Inches(0.4), Inches(1.25), Inches(6.0), Inches(2.5),
                "Standard CoT",
                ["Provide (question, reasoning_steps, answer) examples",
                 "Model learns to show its work before answering",
                 "Arithmetic accuracy: 18% → 57% (PaLM, GSM8K)",
                 "Format: 'Let me think step by step...' as last line of each example"],
                title_color=ACCENT_GREEN, font_size=13)

    bullet_card(slide, Inches(0.4), Inches(3.9), Inches(6.0), Inches(2.5),
                "Zero-Shot CoT",
                ["Simply append: 'Let's think step by step.'",
                 "Works because model has seen this pattern during training",
                 "No examples needed — cheapest CoT variant",
                 "Effective on math, logic, commonsense tasks"],
                title_color=ACCENT_BLUE, font_size=13)

    bullet_card(slide, Inches(6.8), Inches(1.25), Inches(6.1), Inches(2.5),
                "Self-Consistency",
                ["Generate N independent CoT paths (temperature > 0)",
                 "Majority vote on final answers",
                 "Arithmetic: 57% → 78% with k=40 paths",
                 "Tradeoff: k× more expensive"],
                title_color=ACCENT_PURP, font_size=13)

    # Accuracy table
    headers = ["Task", "Standard", "Zero-CoT", "Few-CoT", "Self-Consist"]
    rows = [
        ["Arithmetic",   "18%", "57%", "68%", "78%"],
        ["Commonsense",  "52%", "64%", "83%", "89%"],
        ["Symbolic",     "28%", "63%", "80%", "88%"],
    ]
    x0 = Inches(6.8); y0 = Inches(3.9)
    col_ws = [Inches(1.8), Inches(1.05), Inches(1.05), Inches(1.05), Inches(1.16)]
    row_h = Inches(0.38)
    for ci, (hdr, cw) in enumerate(zip(headers, col_ws)):
        cx = x0 + sum(col_ws[:ci])
        add_rect(slide, cx, y0, cw, row_h, ACCENT_GREEN)
        add_text(slide, hdr, cx+Inches(0.03), y0+Inches(0.05),
                 cw-Inches(0.06), row_h-Inches(0.08), font_size=10, bold=True, color=WHITE)
    for ri, row in enumerate(rows):
        bg = CARD_BG if ri % 2 == 0 else RGBColor(0x16, 0x20, 0x33)
        for ci, (cell, cw) in enumerate(zip(row, col_ws)):
            cx = x0 + sum(col_ws[:ci])
            cy = y0 + (ri+1)*row_h
            add_rect(slide, cx, cy, cw, row_h, bg)
            color = ACCENT_GREEN if ci > 0 and ri == 0 and cell == "78%" else LIGHT_GRAY
            add_text(slide, cell, cx+Inches(0.03), cy+Inches(0.05),
                     cw-Inches(0.06), row_h-Inches(0.08), font_size=11, color=color)


def s_advanced_reasoning(prs):
    slide = blank_slide(prs)
    fill_bg(slide)
    slide_header(slide, "Advanced Reasoning Techniques", accent=ACCENT_PURP)

    cards = [
        ("Least-to-Most", ACCENT_BLUE,
         ["Decompose → solve simplest first → use answer as context",
          "Example: 'Days until March 15?' → Sub-Q1: days left in Feb → Sub-Q2: days in March",
          "Best for: multi-step math, legal reasoning, planning",
          "Key insight: each subproblem answer is injected as context"]),
        ("Skeleton-of-Thought", ACCENT_GREEN,
         ["Step 1: ask model to generate an outline (skeleton)",
          "Step 2: fill each section in parallel",
          "2–3× faster with minimal quality loss",
          "Best for: long-form writing, reports, structured answers"]),
        ("Tree of Thoughts (ToT)", ACCENT_PURP,
         ["Explores multiple reasoning branches (BFS/DFS over thoughts)",
          "Each node = a partial reasoning step",
          "Prune branches with a value/score function",
          "Best for: creative writing, planning, game solving"]),
        ("Graph of Thoughts (GoT)", ACCENT_AMBER,
         ["Generalises ToT: allows cycles and merging",
          "Branch A + Branch B can merge to synthesise insights",
          "Best for: code review loops, iterative document revision",
          "Higher overhead; use only when merging is needed"]),
    ]
    for i, (title, color, bullets) in enumerate(cards):
        col = i % 2; row = i // 2
        bullet_card(slide,
                    Inches(0.4 + col * 6.5),
                    Inches(1.3 + row * 2.9),
                    Inches(6.2), Inches(2.6),
                    title, bullets, title_color=color, font_size=12)


def s_generated_knowledge(prs):
    slide = blank_slide(prs)
    fill_bg(slide)
    slide_header(slide, "Generated Knowledge Prompting", accent=ACCENT_AMBER)

    bullet_card(slide, Inches(0.4), Inches(1.25), Inches(6.0), Inches(5.5),
                "How It Works",
                ["Step 1 — Generate facts:\n  'List 5 facts about penguins relevant to bird classification'",
                 "Step 2 — Use facts to answer:\n  'Using only the facts above, is a penguin a bird?'",
                 "Grounding: reduces hallucination by anchoring to generated context",
                 "No external retrieval needed — fast and cheap",
                 "Limitation: facts are from pre-training; may not reflect latest info",
                 "vs. RAG: RAG fetches real docs; GKP synthesises from memory",
                 "Use cases: knowledge-intensive QA, fact verification, tutoring"],
                title_color=ACCENT_AMBER, font_size=13)

    # Visual flow
    flow = [
        ("User Query",         ACCENT_BLUE),
        ("Generate Facts\n(Pass 1)", ACCENT_GREEN),
        ("Facts + Query\n(Pass 2)", ACCENT_PURP),
        ("Grounded Answer",    ACCENT_AMBER),
    ]
    for i, (label, color) in enumerate(flow):
        bx = Inches(7.1 + i * 1.5)
        add_rect(slide, bx, Inches(2.2), Inches(1.3), Inches(1.0), CARD_BG)
        add_rect(slide, bx, Inches(2.2), Inches(1.3), Inches(0.07), color)
        add_text(slide, label, bx+Inches(0.1), Inches(2.3),
                 Inches(1.1), Inches(0.9), font_size=11, color=WHITE)
        if i < len(flow)-1:
            add_text(slide, "→", bx+Inches(1.3), Inches(2.55),
                     Inches(0.25), Inches(0.4), font_size=18, color=MID_GRAY)

    bullet_card(slide, Inches(6.8), Inches(3.5), Inches(6.1), Inches(3.2),
                "GKP vs RAG vs Fine-tuning",
                ["GKP: fast, no infra, limited to pre-training knowledge",
                 "RAG: fetches real/private docs, updatable, needs vector DB",
                 "Fine-tune: bakes knowledge into weights, fast inference, expensive",
                 "Rule of thumb: GKP for reasoning; RAG for facts; FT for style"],
                title_color=ACCENT_GREEN, font_size=12)


def s_rag(prs):
    slide = blank_slide(prs)
    fill_bg(slide)
    slide_header(slide, "Retrieval-Augmented Generation (RAG)", accent=ACCENT_BLUE)

    # Pipeline steps on left
    steps = [
        ("①  Embed query",         "Same embedding model used to index docs"),
        ("②  Vector search",       "Cosine similarity over chunk index"),
        ("③  Retrieve top-k",      "k=3–5 chunks (512 tokens each typical)"),
        ("④  Stuff context",       "Prepend retrieved text above the question"),
        ("⑤  LLM generates",       "Grounded answer + cites source chunks"),
    ]
    for i, (step, desc) in enumerate(steps):
        y = Inches(1.35 + i * 1.05)
        add_rect(slide, Inches(0.4), y, Inches(6.0), Inches(0.85), CARD_BG)
        add_rect(slide, Inches(0.4), y, Inches(0.06), Inches(0.85), ACCENT_BLUE)
        add_text(slide, step, Inches(0.55), y+Inches(0.05),
                 Inches(5.7), Inches(0.38), font_size=13, bold=True, color=ACCENT_BLUE)
        add_text(slide, desc, Inches(0.55), y+Inches(0.44),
                 Inches(5.7), Inches(0.32), font_size=12, color=LIGHT_GRAY)

    # Chunking strategies on right
    bullet_card(slide, Inches(6.8), Inches(1.25), Inches(6.1), Inches(2.9),
                "Chunking Strategies",
                ["Fixed-size (512 tok, 50 overlap): simple, works well",
                 "Sentence-based (split on .!?): better for Q&A",
                 "Semantic (cluster by embedding): best quality, expensive",
                 "Parent-child: retrieve small chunks, return full doc"],
                title_color=ACCENT_GREEN, font_size=12)

    bullet_card(slide, Inches(6.8), Inches(4.3), Inches(6.1), Inches(2.4),
                "Key Parameters & Tuning",
                ["Chunk size 512–2048: smaller = precise, larger = context-rich",
                 "Top-k 3–5: more = richer but risks irrelevance",
                 "Reranking: cross-encoder re-scores top-20 → keep top-3",
                 "Embedding model choice matters (e5-large, text-embedding-3)"],
                title_color=ACCENT_AMBER, font_size=12)


def s_function_calling(prs):
    slide = blank_slide(prs)
    fill_bg(slide)
    slide_header(slide, "Function Calling & Tool Use", accent=ACCENT_GREEN)

    bullet_card(slide, Inches(0.4), Inches(1.25), Inches(5.7), Inches(2.5),
                "What is it?",
                ["LLM outputs structured JSON call instead of plain text",
                 "System executes the function → returns result → LLM continues",
                 'Output: {"tool": "get_weather", "args": {"city": "London"}}',
                 "Analogy: secretary who knows when to look something up"],
                title_color=ACCENT_GREEN, font_size=13)

    bullet_card(slide, Inches(0.4), Inches(3.85), Inches(5.7), Inches(2.5),
                "When to Use Tool Use",
                ["Real-time data: weather, prices, stock quotes",
                 "Code execution / arithmetic calculations",
                 "Database queries or file I/O",
                 "API calls to external services (calendar, email)",
                 "Any task where you need ground truth, not generation"],
                title_color=ACCENT_BLUE, font_size=13)

    # Tool dispatch loop
    loop_steps = [
        ("Send",      "system + tools + user msg",    ACCENT_BLUE),
        ("Check",     "plain text OR tool_call?",      ACCENT_PURP),
        ("Execute",   "run function, get result",      ACCENT_GREEN),
        ("Append",    "result as tool role message",   ACCENT_AMBER),
        ("Repeat",    "until model returns plain text", ACCENT_ROSE),
    ]
    for i, (title, desc, color) in enumerate(loop_steps):
        y = Inches(1.3 + i * 1.14)
        add_rect(slide, Inches(6.7), y, Inches(6.2), Inches(0.9), CARD_BG)
        add_rect(slide, Inches(6.7), y, Inches(0.06), Inches(0.9), color)
        add_text(slide, f"{i+1}. {title}", Inches(6.85), y+Inches(0.05),
                 Inches(1.3), Inches(0.38), font_size=13, bold=True, color=color)
        add_text(slide, desc, Inches(8.2), y+Inches(0.22),
                 Inches(4.5), Inches(0.45), font_size=12, color=LIGHT_GRAY)


def s_hallucination(prs):
    slide = blank_slide(prs)
    fill_bg(slide)
    slide_header(slide, "Hallucination & Prompt Safety", accent=ACCENT_ROSE)

    bullet_card(slide, Inches(0.4), Inches(1.25), Inches(5.8), Inches(2.5),
                "Types of Hallucination",
                ["Intrinsic: contradicts the provided context",
                 "Extrinsic: goes beyond context (unverifiable claims)",
                 "Factual: incorrect real-world claims (invented citations)",
                 "Root cause: model maximises fluency, not truthfulness",
                 "Frequency: ~20% of medical QA; ~5% on factual benchmarks"],
                title_color=ACCENT_ROSE, font_size=13)

    bullet_card(slide, Inches(0.4), Inches(3.9), Inches(5.8), Inches(2.5),
                "Mitigation Strategies",
                ["1. Retrieval grounding: 'Only use info from context'",
                 "2. Self-check: 'Review above for claims not in context'",
                 "3. Self-consistency: 3/5 paths agree → higher confidence",
                 "4. Calibrated prompt: 'Say I don't know if unsure'",
                 "5. Temperature=0: reduces variability on factual tasks"],
                title_color=ACCENT_GREEN, font_size=13)

    bullet_card(slide, Inches(6.8), Inches(1.25), Inches(6.1), Inches(2.5),
                "Prompt Injection Attacks",
                ["Direct: user message contains 'Ignore all previous instructions'",
                 "Indirect: malicious content in retrieved doc",
                 "Defense: XML-delimit user content:",
                 "  <user_input>{query}</user_input>  (never f-string raw input)",
                 "Validate: strip/flag known injection patterns in preprocessing"],
                title_color=ACCENT_AMBER, font_size=13)

    bullet_card(slide, Inches(6.8), Inches(3.9), Inches(6.1), Inches(2.5),
                "Safe Prompt Patterns",
                ["✓  System prompt sets clear constraints before user content",
                 "✓  Delimiters separate trusted vs. untrusted text",
                 "✓  Constrained output format (JSON schema) limits deviation",
                 "✗  Never interpolate raw user text into system prompts",
                 "✗  Never trust model to enforce security on its own"],
                title_color=ACCENT_PURP, font_size=13)


def s_evaluation(prs):
    slide = blank_slide(prs)
    fill_bg(slide)
    slide_header(slide, "Evaluation & Iteration", accent=ACCENT_AMBER)

    # Metrics table
    headers = ["Metric", "What It Measures", "Range", "Limitation"]
    rows = [
        ["BLEU",         "n-gram precision vs reference",        "0–1", "Penalises paraphrase"],
        ["ROUGE-L",      "Longest common subsequence recall",    "0–1", "No semantics"],
        ["BERTScore",    "Semantic similarity (BERT cosine)",    "0–1", "Slow, expensive"],
        ["Exact Match",  "Binary: output == reference",          "0/1", "Too strict"],
        ["LLM-as-Judge", "Another LLM scores 1–10",              "1–10","Self-bias"],
    ]
    x0 = Inches(0.4); y0 = Inches(1.3)
    col_ws = [Inches(2.0), Inches(3.6), Inches(1.2), Inches(2.5)]
    row_h = Inches(0.42)
    for ci, (hdr, cw) in enumerate(zip(headers, col_ws)):
        cx = x0 + sum(col_ws[:ci])
        add_rect(slide, cx, y0, cw, row_h, ACCENT_AMBER)
        add_text(slide, hdr, cx+Inches(0.04), y0+Inches(0.07),
                 cw-Inches(0.08), row_h-Inches(0.08),
                 font_size=11, bold=True, color=DARK_BG)
    for ri, row in enumerate(rows):
        bg = CARD_BG if ri % 2 == 0 else RGBColor(0x16, 0x20, 0x33)
        for ci, (cell, cw) in enumerate(zip(row, col_ws)):
            cx = x0 + sum(col_ws[:ci])
            cy = y0 + (ri+1)*row_h
            add_rect(slide, cx, cy, cw, row_h, bg)
            add_text(slide, cell, cx+Inches(0.04), cy+Inches(0.07),
                     cw-Inches(0.08), row_h-Inches(0.08), font_size=11, color=LIGHT_GRAY)

    bullet_card(slide, Inches(9.8), Inches(1.25), Inches(3.1), Inches(2.5),
                "A/B Testing Loop",
                ["1. Define success metric",
                 "2. 50–100 test cases",
                 "3. Run Prompt A → score",
                 "4. Run Prompt B → score",
                 "5. Bootstrap significance test",
                 "6. Deploy winner; log loser"],
                title_color=ACCENT_BLUE, font_size=11)

    bullet_card(slide, Inches(9.8), Inches(3.9), Inches(3.1), Inches(2.9),
                "Few-Shot vs. Fine-tune",
                ["<10 examples → zero/few-shot",
                 "10–1K, simple task → few-shot CoT",
                 ">1K, consistent task → LoRA fine-tune",
                 "Latency critical → fine-tune 7B",
                 "Privacy-sensitive → self-hosted + FT"],
                title_color=ACCENT_PURP, font_size=11)


def s_production_checklist(prs):
    slide = blank_slide(prs)
    fill_bg(slide)
    slide_header(slide, "Production Checklist & Best Practices", accent=ACCENT_GREEN)

    items = [
        ("Logging",        ACCENT_BLUE,
         ["Log all prompts, outputs, latency, token count",
          "Tag with user ID and session ID for debugging",
          "Retention: 30-day hot, 1-year cold storage"]),
        ("Monitoring",     ACCENT_GREEN,
         ["Hallucination rate: sample 5% weekly for human review",
          "Refusal rate: track separately from hallucination",
          "Cost/token: alert if spend spikes >20% week-on-week"]),
        ("Versioning",     ACCENT_PURP,
         ["Git-tag prompt configs alongside model version",
          "Maintain regression suite: 100 golden (input, expected) pairs",
          "Reject any change that degrades regression by >2%"]),
        ("Security",       ACCENT_AMBER,
         ["Delimit all user input with XML tags",
          "Block known injection patterns in preprocessing",
          "Rate-limit tool calls per user per minute"]),
        ("Cost Control",   ACCENT_ROSE,
         ["Prompt caching: keep static system prompt prefix identical",
          "Model routing: simple queries → fast/cheap model; hard → flagship",
          "Batch non-latency-sensitive calls for 50% cost saving"]),
        ("Iteration",      ACCENT_BLUE,
         ["Weekly: review sampled outputs with domain experts",
          "Monthly: full regression suite run + benchmark update",
          "Quarterly: consider fine-tuning if >10K good examples accumulated"]),
    ]
    for i, (title, color, bullets) in enumerate(items):
        col = i % 3; row = i // 3
        bullet_card(slide,
                    Inches(0.4 + col * 4.3),
                    Inches(1.3 + row * 2.85),
                    Inches(4.1), Inches(2.6),
                    title, bullets, title_color=color, font_size=12)


def s_key_takeaways(prs):
    slide = blank_slide(prs)
    fill_bg(slide)
    slide_header(slide, "Key Takeaways", accent=ACCENT_PURP)

    points = [
        (ACCENT_BLUE,  "Start simple",
         "Zero-shot first. Add few-shot examples only when zero-shot fails. CoT only for multi-step reasoning."),
        (ACCENT_GREEN, "Ground your prompts",
         "For factual tasks, use RAG or generated knowledge. Never trust the model's memory for critical facts."),
        (ACCENT_PURP,  "Measure everything",
         "Define success metrics before writing prompts. A/B test with ≥50 cases. Track regressions."),
        (ACCENT_AMBER, "Security is not optional",
         "Delimit user input. Validate tool args. Assume adversarial inputs in any public-facing system."),
        (ACCENT_ROSE,  "Iterate systematically",
         "Log → sample → review → update. One change at a time. Git-tag every prompt version."),
        (ACCENT_GREEN, "Cost ≠ quality",
         "Smaller models with good prompts often beat large models with bad ones. Profile before scaling."),
    ]
    for i, (color, title, desc) in enumerate(points):
        col = i % 2; row = i // 2
        x = Inches(0.4 + col * 6.5); y = Inches(1.4 + row * 1.9)
        add_rect(slide, x, y, Inches(6.2), Inches(1.65), CARD_BG)
        add_rect(slide, x, y, Inches(0.08), Inches(1.65), color)
        add_text(slide, title, x+Inches(0.18), y+Inches(0.12),
                 Inches(5.8), Inches(0.4), font_size=15, bold=True, color=color)
        add_text(slide, desc, x+Inches(0.18), y+Inches(0.55),
                 Inches(5.8), Inches(1.0), font_size=13, color=LIGHT_GRAY)


def s_references(prs):
    slide = blank_slide(prs)
    fill_bg(slide)
    slide_header(slide, "References", accent=MID_GRAY)
    refs = [
        "Vaswani et al. (2017). Attention Is All You Need. NeurIPS.",
        "Brown et al. (2020). Language Models are Few-Shot Learners. NeurIPS. [GPT-3]",
        "Wei et al. (2022). Chain-of-Thought Prompting Elicits Reasoning. NeurIPS.",
        "Wang et al. (2023). Self-Consistency Improves CoT Reasoning. ICLR.",
        "Yao et al. (2023). Tree of Thoughts: Deliberate Problem Solving. NeurIPS.",
        "Liu et al. (2022). Generated Knowledge Prompting for Commonsense. ACL.",
        "Lewis et al. (2020). Retrieval-Augmented Generation for NLP. NeurIPS.",
        "Ouyang et al. (2022). Training language models to follow instructions (InstructGPT/RLHF).",
        "Zhou et al. (2023). Least-to-Most Prompting Enables Complex Reasoning. ICLR.",
        "Bai et al. (2022). Constitutional AI: Harmlessness from AI Feedback.",
    ]
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(12.3), Inches(5.7))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, ref in enumerate(refs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_before = Pt(6)
        run = p.add_run()
        run.text = f"[{i+1}]  {ref}"
        run.font.size = Pt(13)
        run.font.color.rgb = LIGHT_GRAY


# ── MAIN ─────────────────────────────────────────────────────────────────────

def main():
    prs = new_prs()

    s_title(prs)
    s_agenda(prs)

    section_divider(prs, "Part 1: Foundations",
                    "LLM internals — architecture, scale, training pipeline")
    s_what_is_llm(prs)
    s_transformer(prs)
    s_training_pipeline(prs)

    section_divider(prs, "Part 2: Prompting Techniques",
                    "From zero-shot to advanced reasoning", color=ACCENT_GREEN)
    s_prompt_anatomy(prs)
    s_nshot(prs)
    s_cot(prs)
    s_advanced_reasoning(prs)
    s_generated_knowledge(prs)

    section_divider(prs, "Part 3: Production Patterns",
                    "RAG, tool use, hallucination, evaluation", color=ACCENT_PURP)
    s_rag(prs)
    s_function_calling(prs)
    s_hallucination(prs)
    s_evaluation(prs)

    section_divider(prs, "Part 4: Best Practices",
                    "Checklists, iteration, and key takeaways", color=ACCENT_AMBER)
    s_production_checklist(prs)
    s_key_takeaways(prs)
    s_references(prs)

    out_path = os.path.join(os.path.dirname(__file__), '..', 'html',
                            'llm_prompt_engineering.pptx')
    out_path = os.path.normpath(out_path)
    prs.save(out_path)
    print(f"Saved {len(prs.slides)} slides to {out_path}")


if __name__ == '__main__':
    main()
